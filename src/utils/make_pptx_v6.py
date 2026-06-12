"""
Build SENTINEL-GNSS Presentation V6.

Changes vs V5:
  - Insert a full-bleed system-pipeline diagram slide at position 8
    (right after the "3 · Method" section divider, before CLEAN/WARN/DEG slide).

Run:
    python -m src.utils.make_pptx_v6
"""

import sys, io, copy, os
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import lxml.etree as etree

SRC = "docs/SENTINEL_GNSS_Presentation_V5.pptx"
DST = "docs/SENTINEL_GNSS_Presentation_V6.pptx"
DIAG = "results/paper_figures/system_pipeline_diagram.png"

NAVY     = RGBColor(0x00, 0x33, 0x66)
BLUE     = RGBColor(0x00, 0x38, 0x93)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
CYAN     = RGBColor(0x4F, 0xC3, 0xF7)

W = Inches(13.33)
H = Inches(7.50)


# ── Helper: copy a blank slide element from the layout ─────────────────────

def blank_slide(prs):
    """Add a new blank slide (layout 6) and return it."""
    blank_layout = prs.slide_layouts[6]   # 'Blank'
    return prs.slides.add_slide(blank_layout)


def move_slide(prs, old_idx: int, new_idx: int):
    """Move slide at old_idx to new_idx (0-based)."""
    xml_slides = prs.slides._sldIdLst
    slides     = list(xml_slides)
    el = slides[old_idx]
    xml_slides.remove(el)
    xml_slides.insert(new_idx, el)


def add_text(slide, text, x, y, w, h,
             size=12, bold=False, color=RGBColor(0x0A, 0x17, 0x33),
             align=PP_ALIGN.LEFT):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = False
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color


# ── Build the pipeline overview slide ──────────────────────────────────────

def build_pipeline_slide(prs, img_path: str):
    """
    Creates a slide with:
      - Navy header bar + title
      - Full-width pipeline diagram image
      - Thin footer note
    """
    slide = blank_slide(prs)
    sp_tree = slide.shapes._spTree

    # --- Header bar (navy, full width, 0.60 in tall) ---
    from pptx.util import Inches, Pt, Emu
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import nsmap

    def solid_rect(slide, x, y, w, h, rgb: RGBColor, zorder=None):
        sp = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            x, y, w, h
        )
        sp.line.fill.background()        # no border
        sp.fill.solid()
        sp.fill.fore_color.rgb = rgb
        sp.line.color.rgb = rgb
        return sp

    # Navy header
    solid_rect(slide, Inches(0), Inches(0), W, Inches(0.58), NAVY)

    # Cyan accent line
    solid_rect(slide, Inches(0), Inches(0.58), W, Inches(0.045), CYAN)

    # Header title
    add_text(slide,
             "COMPLETE SYSTEM PIPELINE  —  SENTINEL-GNSS",
             Inches(0.25), Inches(0.06),
             Inches(10.00), Inches(0.48),
             size=20, bold=True, color=WHITE)

    # Sub-label top-right
    add_text(slide,
             "BEIHANG UNIVERSITY  |  SENTINEL-GNSS · PILOT",
             Inches(8.30), Inches(0.08),
             Inches(4.90), Inches(0.42),
             size=9, bold=False, color=CYAN,
             align=PP_ALIGN.RIGHT)

    # --- Pipeline diagram image ---
    # Position: leave 0.64 in at top (header), 0.32 in at bottom (footer)
    img_y = Inches(0.64)
    img_h = Inches(6.54)   # 7.50 - 0.64 - 0.32
    # Maintain 13.33×6.00 aspect ratio → actual image ratio = 13.33/6.00 = 2.222
    # Fit to slide width (13.33 in) keeping ratio:
    img_w_in = 13.33
    img_h_from_ratio = 13.33 / (13.33 / 6.00)  # = 6.00 in
    # Centre vertically in the available space
    avail = 6.86   # 7.50 - 0.64
    centre_y = Inches(0.64 + (avail - img_h_from_ratio) / 2)

    slide.shapes.add_picture(
        img_path,
        Inches(0), centre_y,
        Inches(img_w_in), Inches(img_h_from_ratio)
    )

    # --- Footer bar ---
    solid_rect(slide, Inches(0), Inches(7.18), W, Inches(0.32), NAVY)
    add_text(slide,
             "SENSORS  →  FEATURE ENGINEERING  →  SENTINEL MODEL  "
             "→  ADAPTIVE EKF  →  FILTERED POSITION  →  DASHBOARD",
             Inches(0.25), Inches(7.19),
             Inches(12.80), Inches(0.30),
             size=8.5, bold=False, color=CYAN)

    return slide


# ── Main ───────────────────────────────────────────────────────────────────

def main():
    assert os.path.exists(SRC),  f"Source not found: {SRC}"
    assert os.path.exists(DIAG), f"Diagram not found: {DIAG}"

    prs = Presentation(SRC)

    n_before = len(prs.slides)
    print(f"Loaded V5: {n_before} slides")

    # Build the pipeline slide (appended at the end for now)
    build_pipeline_slide(prs, DIAG)
    print(f"Pipeline slide added (currently at position {len(prs.slides)})")

    # Move it to position 8 (index 7), right after "3 · Method" divider
    move_slide(prs, len(prs.slides) - 1, 7)
    print("Moved to position 8")

    prs.save(DST)
    n_after = len(prs.slides)
    sz = os.path.getsize(DST) / 1024 / 1024
    print(f"Saved: {DST}  ({n_after} slides, {sz:.1f} MB)")


if __name__ == "__main__":
    main()
