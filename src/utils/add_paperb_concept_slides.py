"""
add_paperb_concept_slides.py
-----------------------------
Build an updated COPY of docs/SENTINEL_GNSS_Presentation.pptx with two new
on-brand slides inserted after the EKF-results slide (slide 25):

  A. "Why raw GPS sometimes beats the EKF" (median vs tail; CEP table)
  B. "Can the noise covariance be computed? Yes" (principled R + online estimator)

The original is left untouched (it may be open in PowerPoint). Output:
  docs/SENTINEL_GNSS_Presentation_updated.pptx

Run: python -m src.utils.add_paperb_concept_slides
"""
from __future__ import annotations
from pathlib import Path
import copy

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).resolve().parents[2]
SRC = ROOT / "docs" / "SENTINEL_GNSS_Presentation.pptx"
OUT = ROOT / "docs" / "SENTINEL_GNSS_Presentation_updated.pptx"

NAVY = RGBColor(0x0A, 0x17, 0x33)
BLUE = RGBColor(0x00, 0x38, 0x93)
GREEN = RGBColor(0x1B, 0x87, 0x3A)
RED = RGBColor(0xC6, 0x28, 0x28)
GREY = RGBColor(0x84, 0x93, 0xA8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF2, 0xF5, 0xFA)
W, H = Inches(13.333), Inches(7.5)


def _blank_layout(prs):
    for lay in prs.slide_layouts:
        if lay.name.strip().lower() == "blank":
            return lay
    return prs.slide_layouts[-1]


def _box(slide, l, t, w, h, fill=None, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
    return sp


def _text(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=4):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (txt, sz, col, bold) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        r = p.add_run(); r.text = txt
        r.font.size = Pt(sz); r.font.bold = bold
        r.font.color.rgb = col; r.font.name = "Calibri"
    return tb


def _chrome(slide, title):
    _box(slide, 0, 0, W, Inches(0.8), fill=NAVY)
    _box(slide, 0, Inches(7.34), W, Inches(0.16), fill=BLUE)
    _text(slide, Inches(0.45), Inches(0.08), Inches(12.4), Inches(0.64),
          [(title, 24, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)


def add_slide(prs, builder):
    slide = prs.slides.add_slide(_blank_layout(prs))
    builder(slide)
    return slide


def slide_median_vs_tail(s):
    _chrome(s, "Why raw GPS sometimes beats the EKF — and why that is correct")
    _text(s, Inches(0.5), Inches(0.95), Inches(12.3), Inches(0.6),
          [("Judge the filter by the worst case, not the typical point. "
            "RMSE is dominated by the tail because it squares errors.",
            15, NAVY, False)])
    # table
    from pptx.util import Inches as I
    rows, cols = 4, 4
    tbl = s.shapes.add_table(rows, cols, I(0.6), I(1.7), I(8.4), I(2.5)).table
    tbl.columns[0].width = I(3.3)
    for c in (1, 2, 3):
        tbl.columns[c].width = I(1.7)
    data = [
        ["Degraded-epoch error", "Raw GPS", "EKF fixed-R", "Winner"],
        ["Median (CEP50)", "6.7 m", "9.7 m", "Raw GPS"],
        ["95th pct (CEP95)", "76 m", "51 m", "EKF"],
        ["Worst case (max)", "888 m", "137 m", "EKF"],
    ]
    win_col = {"Raw GPS": GREY, "EKF": GREEN}
    for ri in range(rows):
        for ci in range(cols):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if ri == 0 else (LIGHT if ri % 2 else WHITE)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            r = p.add_run(); r.text = data[ri][ci]
            r.font.size = Pt(14 if ri == 0 else 13)
            r.font.name = "Calibri"
            r.font.bold = (ri == 0 or ci == 3)
            if ri == 0:
                r.font.color.rgb = WHITE
            elif ci == 3:
                r.font.color.rgb = win_col.get(data[ri][ci], NAVY)
            else:
                r.font.color.rgb = NAVY
    # right rail: the takeaway
    _box(s, Inches(9.25), Inches(1.7), Inches(3.55), Inches(2.5), fill=LIGHT, line=BLUE)
    _text(s, Inches(9.45), Inches(1.85), Inches(3.2), Inches(2.25),
          [("The trade you want", 15, BLUE, True),
           ("Raw GPS wins the median — the harmless points.", 12.5, NAVY, False),
           ("The EKF wins the tail — the dangerous points "
            "(no 888 m excursions).", 12.5, NAVY, False),
           ("Give up ~2 m where it doesn't matter, to never be "
            "800 m wrong where it does.", 12.5, GREEN, True)], space=8)
    # bottom message bar
    _box(s, Inches(0.6), Inches(4.45), Inches(12.2), Inches(2.5), fill=NAVY)
    _text(s, Inches(0.9), Inches(4.65), Inches(11.6), Inches(2.1),
          [("Goal: when GPS degrades, the vehicle must still know where it is well "
            "enough to act.", 16, WHITE, True),
           ("A position estimate that is usually within 7 m but sometimes 888 m off is "
            "far more dangerous than one that is usually within 10 m and never worse "
            "than 137 m.", 14, RGBColor(0xCF, 0xD8, 0xE8), False),
           ("The visible offsets where GPS is good are filter lag + deliberate distrust — "
            "the unavoidable price of not lurching onto BAD GPS in the canyon (the "
            "“swirling” failure).", 14, RGBColor(0xCF, 0xD8, 0xE8), False)], space=10)


def slide_principled_r(s):
    _chrome(s, "Can the noise covariance R be computed?  Yes.")
    _text(s, Inches(0.5), Inches(0.95), Inches(12.3), Inches(0.5),
          [("A fixed r_base is only a baseline. Principled, per-epoch R is available "
            "(in rough order of rigour):", 15, NAVY, False)])
    methods = [
        ("1. Receiver-reported sigma", "R = (sigma_reported)² per epoch — already in our pipeline (spp_horiz_std)."),
        ("2. Geometry-based", "R ∝ (HDOP · UERE)² — turns satellite geometry into a position variance."),
        ("3. C/N₀-weighted", "Per-satellite variance as a function of signal strength."),
        ("4. Innovation-based (Mehra)", "Principled in theory, but diverges in urban NLOS (238,000 m) — disabled."),
        ("5. Online σ_deg estimator (adopted)", "Causal running median of degraded-epoch innovations."),
    ]
    y = 1.55
    for head, body in methods:
        col = GREEN if head.startswith("5") else (RED if head.startswith("4") else BLUE)
        _box(s, Inches(0.6), Inches(y), Inches(0.12), Inches(0.66), fill=col)
        _text(s, Inches(0.85), Inches(y), Inches(11.9), Inches(0.66),
              [(head + "  —  " + body, 13.5, NAVY, head.startswith("5"))],
              anchor=MSO_ANCHOR.MIDDLE)
        y += 0.78
    # result banner
    _box(s, Inches(0.6), Inches(5.6), Inches(12.2), Inches(1.35), fill=NAVY)
    _text(s, Inches(0.9), Inches(5.72), Inches(11.6), Inches(1.15),
          [("Online estimator removes per-environment hand-tuning AND improves accuracy",
            16, WHITE, True),
           ("Tokyo Shinjuku: calibrated SENTINEL 38.7%  →  + online σ_deg  43.2% "
            "degraded-RMSE reduction (stable across the gate range).",
            14, RGBColor(0xCF, 0xD8, 0xE8), False)], space=8)


def main():
    prs = Presentation(str(SRC))
    n_before = len(prs.slides._sldIdLst)
    add_slide(prs, slide_median_vs_tail)
    add_slide(prs, slide_principled_r)

    # Move the two appended slides to position right after slide 25 (index 24 -> insert at 25,26)
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    new1, new2 = ids[-2], ids[-1]
    sldIdLst.remove(new1); sldIdLst.remove(new2)
    sldIdLst.insert(25, new1)   # becomes slide 26
    sldIdLst.insert(26, new2)   # becomes slide 27

    prs.save(str(OUT))
    print(f"slides before: {n_before}; after: {len(prs.slides._sldIdLst)}")
    print(f"inserted concept slides as #26 and #27 -> {OUT.name}")


if __name__ == "__main__":
    main()
