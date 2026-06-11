"""
Build SENTINEL-GNSS Presentation V5.

Changes vs V4 (June 2026 presentation):
  - Delete blank slide 16 (empty cross-city divider)
  - Fix slide 22 divider: "5.5" -> "5 · EKF: From Prediction to Position"
  - Replace slide 23 (blank) -> EKF Concept (Why the EKF?)
  - Replace slide 24 (blank) -> EKF Adaptive-R Formula
  - Insert new slide after 24: EKF Results (3-tier table)
  - Insert new slide: EKF Severity Sweep (when does adaptive-R help?)
  - Insert new slide: Dashboard Overview
  - Insert new slide: Dashboard Live Demo
  - Update slide 6 (Dataset): add role note (HK=train, Tokyo=held-out)
  - Trim slide 25 (Validation): keep only E3/E4/E6/E7

Run:
    python -m src.utils.make_pptx_v5
"""

import lxml.etree as etree
from pptx.oxml.ns import qn
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt, Emu
from pptx import Presentation
import copy
import sys
import io
import os
sys.stdout = io.TextIOWrapper(
    sys.stdout.buffer, encoding='utf-8', errors='replace')


# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
SRC = ("docs/FOURTH MEETING GROUP3 TOPIC"
       "(AI-Based Prediction of GNSS Signal Degradation for Autonomous Driving)"
       " PRESENTED ON 2026-06-04.pptx")
DST = ("docs/SENTINEL_GNSS_Presentation_V5.pptx")

# ---------------------------------------------------------------------------
# Colour palette (match existing design)
# ---------------------------------------------------------------------------
NAVY = RGBColor(0x00, 0x33, 0x66)   # section divider bg
BLUE_DARK = RGBColor(0x00, 0x38, 0x93)   # headings
BLUE_MID = RGBColor(0x00, 0x5B, 0xAC)   # sub-headings
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x0A, 0x17, 0x33)
GREY = RGBColor(0x33, 0x33, 0x33)
GREY2 = RGBColor(0x55, 0x55, 0x55)
CYAN = RGBColor(0x4F, 0xC3, 0xF7)   # callout accent
GREEN = RGBColor(0x1B, 0x87, 0x3A)
AMBER = RGBColor(0xF5, 0x7F, 0x17)
RED = RGBColor(0xC6, 0x28, 0x28)

W = Inches(13.33)
H = Inches(7.50)

# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------


def _sp(slide):
    """Return the spTree element of the slide."""
    return slide.shapes._spTree


def add_box(slide, text, left, top, width, height,
            size=18, color=DARK, bold=False,
            align=PP_ALIGN.LEFT, wrap=True, italic=False):
    """Add a plain text box."""
    from pptx.util import Pt
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return txBox


def add_box_ml(slide, lines, left, top, width, height,
               size=16, bold_first=False):
    """
    Add multi-line text box.
    lines: list of (text, size, color, bold)
    """
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for (txt, sz, col, bld) in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(sz)
        run.font.color.rgb = col
        run.font.bold = bld
    return txBox


def add_filled_box(slide, left, top, width, height, fill_color, alpha=None):
    """Add a solid-filled rectangle (no text)."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_table_shape(slide, rows, cols, left, top, width, height,
                    header_data, body_data,
                    hdr_bg=NAVY, hdr_fg=WHITE,
                    alt_bg=RGBColor(0xE8, 0xF0, 0xFE),
                    row_bg=WHITE):
    """Add a formatted table."""
    tbl = slide.shapes.add_table(
        rows, cols, Inches(left), Inches(top),
        Inches(width), Inches(height)).table

    # Column widths (equal split)
    col_w = Inches(width / cols)
    for c in range(cols):
        tbl.columns[c].width = col_w

    # Header row
    for c, hdr in enumerate(header_data):
        cell = tbl.cell(0, c)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = hdr_bg
        tf = cell.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].runs[0] if tf.paragraphs[0].runs else tf.paragraphs[0].add_run(
        )
        run.text = hdr
        run.font.bold = True
        run.font.size = Pt(12)
        run.font.color.rgb = hdr_fg

    # Body rows
    for r, row_data in enumerate(body_data):
        bg = alt_bg if r % 2 == 0 else row_bg
        for c, val in enumerate(row_data):
            cell = tbl.cell(r + 1, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            tf = cell.text_frame
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            run = tf.paragraphs[0].runs[0] if tf.paragraphs[0].runs else tf.paragraphs[0].add_run(
            )
            run.text = val
            run.font.size = Pt(11)
            run.font.color.rgb = DARK
            if val.startswith('+') or val.startswith('−') or 'wins' in val.lower() or '%' in val:
                run.font.bold = True
                if val.startswith('+') or 'wins' in val.lower():
                    run.font.color.rgb = GREEN
                elif '−' in val or val.startswith('-'):
                    run.font.color.rgb = RED

    return tbl


# ---------------------------------------------------------------------------
# Slide frame (header bar + logo text) -- applied to blank new slides
# ---------------------------------------------------------------------------

def add_frame(slide, title_text):
    """Add the top bar, bottom bar, and title label matching the presentation style."""
    # Top bar (dark navy)
    bar = slide.shapes.add_shape(1,
                                 Inches(0), Inches(-0.03), Inches(13.33), Inches(0.84))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    # Top accent line (cyan)
    acc = slide.shapes.add_shape(1,
                                 Inches(0), Inches(0.77), Inches(13.33), Inches(0.04))
    acc.fill.solid()
    acc.fill.fore_color.rgb = CYAN
    acc.line.fill.background()

    # Bottom bar
    bot = slide.shapes.add_shape(1,
                                 Inches(0), Inches(7.27), Inches(13.33), Inches(0.23))
    bot.fill.solid()
    bot.fill.fore_color.rgb = NAVY
    bot.line.fill.background()

    # Bottom accent
    bacc = slide.shapes.add_shape(1,
                                  Inches(0), Inches(7.24), Inches(13.33), Inches(0.03))
    bacc.fill.solid()
    bacc.fill.fore_color.rgb = CYAN
    bacc.line.fill.background()

    # Title label in top bar
    add_box(slide, title_text,
            left=1.39, top=0.07, width=11.5, height=0.66,
            size=18, color=WHITE, bold=True, align=PP_ALIGN.LEFT)

    # Footer text
    add_box(slide, "BEIHANG UNIVERSITY",
            left=0.67, top=7.30, width=4.0, height=0.17,
            size=8, color=RGBColor(0xBF, 0xE6, 0xFB), bold=False)
    add_box(slide, "SENTINEL-GNSS · PILOT",
            left=9.5, top=7.30, width=3.2, height=0.17,
            size=8, color=WHITE, bold=False, align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------------------
# Section divider slide (navy background, big white text)
# ---------------------------------------------------------------------------

def make_divider_slide(prs, text):
    """Create a section-divider slide (solid navy, large white centred text)."""
    blank_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(blank_layout)

    # Solid navy background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = NAVY

    # Big white centred text
    add_box(slide, text,
            left=0.90, top=2.80, width=11.50, height=1.20,
            size=40, color=WHITE, bold=True, align=PP_ALIGN.LEFT)

    # Accent underline
    acc = slide.shapes.add_shape(1,
                                 Inches(0.95), Inches(4.15), Inches(3.20), Inches(0.07))
    acc.fill.solid()
    acc.fill.fore_color.rgb = CYAN
    acc.line.fill.background()

    return slide


# ---------------------------------------------------------------------------
# Slide deletion / insertion helpers
# ---------------------------------------------------------------------------

def delete_slide(prs, index):
    """Delete slide at 0-based index, cleanly removing it from the package."""
    xml_slides = prs.slides._sldIdLst
    slide = prs.slides[index]
    slide_part = slide.part
    prs_part = prs.part

    # Find the relationship ID
    rId = None
    for rel in prs_part.rels.values():
        if rel._target is slide_part:
            rId = rel.rId
            break

    # Remove from sldIdLst
    xml_slides.remove(xml_slides[index])

    # Drop the relationship from the presentation part
    if rId:
        prs_part.drop_rel(rId)

    # Remove the slide part from the package's _parts registry so the
    # partname is freed and won't produce a duplicate-name ZIP warning.
    try:
        pkg = slide_part._package
        partname = slide_part.partname
        if partname in pkg._parts:
            del pkg._parts[partname]
        # Also drop the slide's own rels part if present
        rels_partname = partname.baseURI + '/_rels/' + partname.filename + '.rels'
        if rels_partname in pkg._parts:
            del pkg._parts[rels_partname]
    except Exception as e:
        print(f"  (package cleanup warning: {e})")


def move_slide(prs, old_index, new_index):
    """Move slide from old_index to new_index (0-based)."""
    xml_slides = prs.slides._sldIdLst
    slide = xml_slides[old_index]
    xml_slides.remove(slide)
    xml_slides.insert(new_index, slide)


# ---------------------------------------------------------------------------
# Build the new slides
# ---------------------------------------------------------------------------

def build_ekf_why(prs):
    """EKF Slide 1: Why the Adaptive EKF?"""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_frame(slide, "WHY THE ADAPTIVE EKF?")

    # Subtitle
    add_box(slide, "Prediction is only valuable if it improves actual navigation",
            left=0.67, top=0.98, width=12.0, height=0.40,
            size=14, color=GREY, bold=False, align=PP_ALIGN.LEFT)

    # Left column header
    add_filled_box(slide, 0.67, 1.55, 5.90, 0.42,
                   RGBColor(0xC6, 0x28, 0x28))
    add_box(slide, "Standard Kalman Filter",
            left=0.67, top=1.55, width=5.90, height=0.42,
            size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Right column header
    add_filled_box(slide, 6.77, 1.55, 5.90, 0.42, GREEN)
    add_box(slide, "Our Adaptive EKF",
            left=6.77, top=1.55, width=5.90, height=0.42,
            size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Left column items
    left_items = [
        ("R = fixed (e.g. 9 m²)", RED),
        ("Always trusts GNSS equally", GREY),
        ("Waits until GNSS fails to react", GREY),
        ("Position jumps during blockage", GREY),
        ("No advance warning used", GREY),
    ]
    top = 2.10
    for txt, col in left_items:
        add_box(slide, "●  " + txt,
                left=0.77, top=top, width=5.70, height=0.40,
                size=13, color=col, bold=(col == RED))
        top += 0.47

    # Right column items
    right_items = [
        ("R(t) grows with P̂(DEGRADED)", GREEN),
        ("Pre-emptively distrusts GNSS", DARK),
        ("Shifts to dead-reckoning before failure", DARK),
        ("Smooth handoff during blockage", DARK),
        ("5 s early warning from SENTINEL", DARK),
    ]
    top = 2.10
    for txt, col in right_items:
        add_box(slide, "✔  " + txt,
                left=6.87, top=top, width=5.70, height=0.40,
                size=13, color=col, bold=(col == GREEN))
        top += 0.47

    # Divider line between columns
    div = slide.shapes.add_shape(1,
                                 Inches(6.60), Inches(1.55), Inches(0.05), Inches(3.00))
    div.fill.solid()
    div.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    div.line.fill.background()

    # Bottom callout bar
    add_filled_box(slide, 0.67, 6.45, 12.0, 0.58,
                   RGBColor(0x00, 0x38, 0x93))
    add_box(slide,
            "“Prediction closes the loop: we don’t wait for GNSS to fail — we pre-empt it”",
            left=0.77, top=6.45, width=11.80, height=0.58,
            size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    return slide


def build_ekf_formula(prs):
    """EKF Slide 2: The Adaptive-R Formula."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_frame(slide, "ADAPTIVE MEASUREMENT NOISE — THE MECHANISM")

    # Main formula box
    add_filled_box(slide, 1.00, 1.05, 11.33, 0.80, RGBColor(0xF0, 0xF4, 0xFF))
    add_box(slide, "R(t)  =  σ²_base  +  (σ²_deg − σ²_base)  ×  P̂_calib(t)",
            left=1.00, top=1.05, width=11.33, height=0.80,
            size=24, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

    # Calibration sub-formula
    add_filled_box(slide, 2.00, 1.95, 9.33, 0.55, RGBColor(0xE8, 0xF0, 0xFE))
    add_box(slide, "P̂_calib(t)  =  clip ( ( P̂(t) − P₅ ) / ( 1 − P₅ ),   0,   1 )",
            left=2.00, top=1.95, width=9.33, height=0.55,
            size=17, color=BLUE_DARK, bold=False, align=PP_ALIGN.CENTER)

    # Kalman gain
    add_box(slide, "Kalman gain:   Kₜ  =  P⁻ₜ Hᵀ ( H P⁻ₜ Hᵀ + Rₜ )⁻¹",
            left=1.00, top=2.62, width=11.33, height=0.45,
            size=15, color=GREY, bold=False, align=PP_ALIGN.CENTER)

    # Annotation arrows label section
    add_box(slide, "What each term means:",
            left=0.67, top=3.18, width=4.0, height=0.30,
            size=12, color=BLUE_MID, bold=True)

    annotations = [
        ("σ²_base  =  9 m²",
         "Baseline GNSS noise (signal is CLEAN) — filter trusts GNSS tightly"),
        ("σ²_deg  =  10,000 m²",
         "Noise under full degradation — filter ignores GNSS, dead-reckons"),
        ("P̂_calib(t)",
         "Calibrated probability from SENTINEL (0 = clean, 1 = fully degraded)"),
        ("P₅  =  0.153",
         "Floor offset — unsupervised calibration removes receiver-domain bias"),
        ("Kₜ  shrinks when Rₜ rises",
         "When SENTINEL predicts degradation, filter leans on motion model"),
    ]
    top = 3.55
    for term, desc in annotations:
        add_box(slide, term,
                left=0.77, top=top, width=2.50, height=0.35,
                size=11, color=BLUE_DARK, bold=True)
        add_box(slide, desc,
                left=3.40, top=top, width=9.30, height=0.35,
                size=11, color=GREY)
        top += 0.42

    # Three value pills at bottom
    pills = [
        (0.67,  "P̂ = 0",   "R = 9 m²  →  Trust GNSS",      GREEN),
        (4.72,  "P̂ = 0.5", "R ≈ 500 m²  →  Caution",   AMBER),
        (8.77,  "P̂ = 1",   "R = 10,000 m²  →  Dead-reckon", RED),
    ]
    for x, label, desc, col in pills:
        add_filled_box(slide, x, 6.30, 3.55, 0.78, col)
        add_box(slide, label,
                left=x, top=6.30, width=3.55, height=0.35,
                size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_box(slide, desc,
                left=x, top=6.62, width=3.55, height=0.35,
                size=10, color=WHITE, bold=False, align=PP_ALIGN.CENTER)

    return slide


def build_ekf_results(prs):
    """EKF Slide 3: Results — Three Tiers."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_frame(slide, "EKF RESULTS — THREE TIERS OF VALIDATION")

    add_box(slide, "From controlled simulation to fully real Tokyo data (cm-level SPAN-INS ground truth)",
            left=0.67, top=0.98, width=12.0, height=0.35,
            size=13, color=GREY, align=PP_ALIGN.LEFT)

    # Big headline numbers
    stats = [
        (0.67,  "−33.8%", "Synthetic blockage\n54.4 m → 36.0 m", RED),
        (4.72,  "+82%",        "Semi-synthetic Tokyo\n36.3 m → 6.4 m", BLUE_DARK),
        (8.77,  "+48.8%",      "Fully real Tokyo\n47.4 m → 24.3 m", GREEN),
    ]
    for x, pct, desc, col in stats:
        add_filled_box(slide, x, 1.45, 3.55, 1.55, RGBColor(0xF5, 0xF7, 0xFF))
        add_box(slide, pct,
                left=x + 0.05, top=1.45, width=3.45, height=0.85,
                size=40, color=col, bold=True, align=PP_ALIGN.CENTER)
        add_box(slide, desc,
                left=x + 0.05, top=2.25, width=3.45, height=0.70,
                size=11, color=DARK, bold=False, align=PP_ALIGN.CENTER)

    # Table
    headers = ["Tier", "Data", "Blocked RMSE", "Gain"]
    body = [
        ["Synthetic",       "Controlled blockage simulation",
            "54.4 m → 36.0 m", "−33.8%"],
        ["Semi-synthetic",  "Real Tokyo path + IMU, synthetic GNSS errors",
            "36.3 m → 6.4 m",  "+82%"],
        ["Fully real ★", "RTKLIB Trimble + real IMU + SPAN-INS truth",
            "47.4 m → 24.3 m", "+48.8%"],
    ]
    add_table_shape(slide, rows=4, cols=4,
                    left=0.67, top=3.15, width=12.0, height=1.70,
                    header_data=headers, body_data=body)

    # Callout
    add_filled_box(slide, 0.67, 5.05, 12.0, 0.70, RGBColor(0xE8, 0xF5, 0xE9))
    add_box(slide,
            "★ Fully real validation: real GNSS + IMU + cm-level ground truth on Tokyo Shinjuku — "
            "the city SENTINEL never saw in training. 47.4 m → 24.3 m blocked-segment RMSE.",
            left=0.77, top=5.05, width=11.80, height=0.70,
            size=12, color=GREEN, bold=True)

    # Honest note
    add_box(slide,
            "Honest note: the dominant win is IMU aiding (odometry + NHC + ZUPT). "
            "Adaptive-R adds on top in severe multipath — quantified in the next slide.",
            left=0.67, top=5.85, width=12.0, height=0.40,
            size=11, color=GREY2, italic=True)

    return slide


def build_ekf_severity(prs):
    """EKF Slide 4: When does adaptive-R help? (Severity Sweep)."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_frame(slide, "WHEN IS ADAPTIVE-R WORTH IT? — SEVERITY SWEEP")

    add_box(slide, "We swept multipath severity to avoid cherry-picking one scenario",
            left=0.67, top=0.98, width=12.0, height=0.35,
            size=13, color=GREY)

    # Left panel — GNSS-only platform
    add_filled_box(slide, 0.67, 1.45, 5.80, 0.42, RED)
    add_box(slide, "GNSS-Only Platform (no IMU, no odometry)",
            left=0.67, top=1.45, width=5.80, height=0.42,
            size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    gnss_only = [
        ("Crossover at ~20 m multipath severity",          RED,      True),
        ("Below 20 m: fixed-R wins (GNSS is fine)",        GREY,     False),
        ("Above 20 m: adaptive-R wins by +25–38%",    GREEN,    True),
        ("Deep canyons, tunnels: exactly our target",      DARK,     False),
        ("Cheap receiver + prediction = strong combo",     DARK,     False),
    ]
    top = 2.00
    for txt, col, bld in gnss_only:
        add_box(slide, "●  " + txt,
                left=0.77, top=top, width=5.60, height=0.38,
                size=12, color=col, bold=bld)
        top += 0.43

    # Right panel — Well-aided platform
    add_filled_box(slide, 6.77, 1.45, 5.89, 0.42, BLUE_DARK)
    add_box(slide, "Well-Aided Platform (odometry + NHC + ZUPT)",
            left=6.77, top=1.45, width=5.89, height=0.42,
            size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    aided = [
        ("Fixed-R wins across realistic range",            GREEN,    True),
        ("GNSS provides the only heading reference",       DARK,     False),
        ("Blanket R-inflation → heading drift",       RED,      True),
        ("Odometry = speed;  NHC = no lateral slip",      GREY,     False),
        ("But heading has no absolute backup",             GREY,     False),
    ]
    top = 2.00
    for txt, col, bld in aided:
        add_box(slide, "●  " + txt,
                left=6.87, top=top, width=5.69, height=0.38,
                size=12, color=col, bold=bld)
        top += 0.43

    # Divider
    div = slide.shapes.add_shape(1,
                                 Inches(6.60), Inches(1.45), Inches(0.05), Inches(2.70))
    div.fill.solid()
    div.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    div.line.fill.background()

    # Practical takeaway
    add_filled_box(slide, 0.67, 4.58, 12.0, 0.55, RGBColor(0xE3, 0xF2, 0xFD))
    add_box(slide,
            "Practical rule: use adaptive-R on GNSS-only platforms in deep-canyon / tunnel environments. "
            "On full AV sensor suite, SENTINEL’s role is integrity flagging and regime selection.",
            left=0.77, top=4.58, width=11.80, height=0.55,
            size=12, color=BLUE_DARK, bold=True)

    # Summary box
    add_box(slide, "Role of SENTINEL in a full AV stack:",
            left=0.67, top=5.25, width=5.50, height=0.30,
            size=12, color=BLUE_MID, bold=True)
    roles = [
        "✔  Trigger sensor-fusion mode switch (GNSS-primary → odometry-primary)",
        "✔  Alert path planner: predicted blockage zone ahead",
        "✔  Activate ZUPT when vehicle is likely stopped in blocked zone",
        "✔  Flag integrity failure to downstream systems (HD map localisation)",
    ]
    top = 5.58
    for r in roles:
        add_box(slide, r, left=0.77, top=top, width=11.80, height=0.33,
                size=11, color=DARK)
        top += 0.35

    return slide


def build_dashboard_overview(prs):
    """Dashboard Slide 1: SENTINEL Dashboard Overview."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_frame(slide, "SENTINEL-GNSS DASHBOARD — REAL-TIME ANALYTICS")

    add_box(slide, "FastAPI backend + Next.js frontend · WebSocket streaming · Runs on any laptop",
            left=0.67, top=0.98, width=12.0, height=0.35,
            size=13, color=GREY, align=PP_ALIGN.LEFT)

    panels = [
        ("01", "Signal Gauge",
         "P(DEGRADED) at +5 / +15 / +30 s\nColour-coded green / amber / red",
         RGBColor(0x1B, 0x87, 0x3A)),
        ("02", "Class Probability Bars",
         "CLEAN / WARNING / DEGRADED\nLive confidence per horizon",
         BLUE_DARK),
        ("03", "Trajectory Map",
         "Vehicle path coloured by\npredicted risk level",
         RGBColor(0x6A, 0x1B, 0x9A)),
        ("04", "P(DEGRADED) Timeline",
         "All 3 horizons streaming\nwith threshold lines",
         RGBColor(0x00, 0x83, 0x8F)),
        ("05", "EKF Analytics",
         "Blocked-segment RMSE by filter\nAided EKF wins at 24.3 m",
         AMBER),
        ("06", "Alert Centre",
         "CRITICAL: P > 0.8 @ +5 s\nWARNING: P > 0.6 @ +15 s",
         RED),
    ]

    cols = 3
    cell_w, cell_h = 4.00, 1.85
    start_x, start_y = 0.67, 1.45
    gap_x, gap_y = 0.17, 0.22

    for i, (num, title, body, col) in enumerate(panels):
        row = i // cols
        col_idx = i % cols
        x = start_x + col_idx * (cell_w + gap_x)
        y = start_y + row * (cell_h + gap_y)

        add_filled_box(slide, x, y, cell_w, cell_h,
                       RGBColor(0xF5, 0xF7, 0xFF))
        # Colored top strip
        add_filled_box(slide, x, y, cell_w, 0.10, col)
        # Number badge
        add_filled_box(slide, x, y + 0.12, 0.38, 0.38, col)
        add_box(slide, num,
                left=x, top=y + 0.12, width=0.38, height=0.38,
                size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_box(slide, title,
                left=x + 0.45, top=y + 0.12, width=cell_w - 0.55, height=0.40,
                size=12, color=DARK, bold=True)
        add_box(slide, body,
                left=x + 0.12, top=y + 0.60, width=cell_w - 0.20, height=0.95,
                size=11, color=GREY)

    # Bottom bar
    add_filled_box(slide, 0.67, 6.88, 12.0, 0.40, NAVY)
    add_box(slide,
            "Pure-SVG visualisations · Zero external dependencies · Works offline · WebSocket @ 1 Hz",
            left=0.77, top=6.88, width=11.80, height=0.40,
            size=12, color=WHITE, bold=False, align=PP_ALIGN.CENTER)

    return slide


def build_dashboard_demo(prs):
    """Dashboard Slide 2: Live Demo — 3 Minutes."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_frame(slide, "LIVE DEMO — 3 MINUTES")

    add_box(slide, "Scenario A — instant blockage · real Beihang campus NMEA · real inference output",
            left=0.67, top=0.98, width=12.0, height=0.35,
            size=13, color=GREY)

    headers = ["Step", "Action", "What the audience sees"]
    steps = [
        ["1", "Open http://localhost:3000",
         "Full dashboard loads — 6 live panels"],
        ["2", "Select “A_log_0000” (instant blockage)",
         "Prediction data populates all panels"],
        ["3", "Press ▶ Play at 5× speed",
         "P(DEGRADED) timeline starts streaming"],
        ["4", "Watch gauge spike before GNSS drops",
         "Gauge turns RED — CRITICAL alert fires"],
        ["5", "Pause — point to lead-time number",
         "“83 m of reaction distance at 60 km/h”"],
        ["6", "Switch to EKF Analytics tab",
         "Blocked-segment RMSE chart: aided EKF wins"],
        ["7", "Point to trajectory map",
         "Path colour shifts green → red through blockage zone"],
    ]

    add_table_shape(slide, rows=8, cols=3,
                    left=0.67, top=1.45, width=12.0, height=4.50,
                    header_data=headers, body_data=steps,
                    hdr_bg=NAVY, hdr_fg=WHITE)

    # Callout
    add_filled_box(slide, 0.67, 6.10, 12.0, 0.60, RGBColor(0xE8, 0xF5, 0xE9))
    add_box(slide,
            "★  Everything is REAL: pre-computed inference on real Beihang NMEA data. "
            "Not a mock-up. Not a demo mode.",
            left=0.77, top=6.10, width=11.80, height=0.60,
            size=13, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

    return slide


# ---------------------------------------------------------------------------
# Dataset slide update: add role note
# ---------------------------------------------------------------------------

def update_dataset_slide(slide):
    """Add a role-clarification note to the dataset slide (slide index 5)."""
    add_filled_box(slide, 0.67, 6.40, 12.0, 0.48, RGBColor(0xFF, 0xF8, 0xE1))
    add_box(slide,
            "Training: Beihang (Hangzhou) field Scenarios A–E  +  UrbanNav Hong Kong (Medium / Deep / Harsh / Tunnel)  —  62,413 windows  |  "
            "Zero-shot test: UrbanNav Tokyo Shinjuku (NEVER in training)",
            left=0.77, top=6.40, width=11.80, height=0.48,
            size=11, color=RGBColor(0x7B, 0x52, 0x00), bold=True)


# ---------------------------------------------------------------------------
# Validation slide trim: rewrite to show only E3/E4/E6/E7
# ---------------------------------------------------------------------------

def rebuild_validation_slide(slide):
    """Clear text boxes on slide 25 and replace with trimmed E3/E4/E6/E7 content."""
    # Remove all existing text-frame shapes except the structural ones (Image 0, bars)
    to_remove = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name.startswith("Text"):
            to_remove.append(shape)
    spTree = slide.shapes._spTree
    for shape in to_remove:
        spTree.remove(shape._element)

    # Rewrite with 4 key experiments
    add_box(slide, "VALIDATION EXPERIMENTS",
            left=1.39, top=0.07, width=11.5, height=0.66,
            size=18, color=WHITE, bold=True)

    experiments = [
        ("E3", "Bootstrap 95% CI (1,000 iterations)",
         "DEGRADED F1 = 0.718 [0.671, 0.762] — all claims backed by confidence intervals. "
         "MCC = 0.773 ± 0.035.",
         BLUE_DARK),
        ("E4", "Inference Latency (Tesla T4 GPU)",
         "SENTINEL-GNSS: 0.039 ms/sample vs 3 tree models: 0.409 ms — 10.52× speedup. "
         "17.8 MB checkpoint. <0.04% of CPU budget at 10 Hz.",
         GREEN),
        ("E6", "Cross-City Zero-Shot — Tokyo Shinjuku",
         "XGBoost DEGRADED F1: 0.822 → 0.148 (−82%). SENTINEL-GNSS: 0.718 → 0.753 (+5%). "
         "DL + XGBoost ensemble: 0.892 macro-F1, 0.896 DEGRADED F1. KEY RESULT.",
         RED),
        ("E7", "Probability Calibration (ECE)",
         "Temperature scaling (T = 0.4023) cuts ECE from 0.114 → 0.068 (−40%). "
         "Reliable P(DEGRADED) enables the adaptive EKF.",
         AMBER),
    ]

    top = 1.10
    for code, title, body, col in experiments:
        add_filled_box(slide, 0.67, top, 0.60, 1.35, col)
        add_box(slide, code,
                left=0.67, top=top + 0.40, width=0.60, height=0.50,
                size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_filled_box(slide, 1.37, top, 11.30, 1.35,
                       RGBColor(0xF5, 0xF7, 0xFF))
        add_box(slide, title,
                left=1.47, top=top + 0.05, width=11.10, height=0.38,
                size=13, color=col, bold=True)
        add_box(slide, body,
                left=1.47, top=top + 0.48, width=11.10, height=0.75,
                size=11, color=GREY)
        top += 1.50


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    prs = Presentation(SRC)
    n = len(prs.slides)
    print(f"Loaded: {n} slides")

    # -----------------------------------------------------------------------
    # 1. Update Dataset slide (index 5)
    # -----------------------------------------------------------------------
    update_dataset_slide(prs.slides[5])
    print("1. Dataset slide updated")

    # -----------------------------------------------------------------------
    # 2. Fix section divider slide 22 (index 21): "5.5" -> "5"
    # -----------------------------------------------------------------------
    div_slide = prs.slides[21]
    for shape in div_slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if "5.5" in run.text or "5.5" in shape.text_frame.text:
                        run.text = run.text.replace(
                            "5.5 · EKF: From Prediction to Position",
                            "5 · EKF: From Prediction to Position")
                        run.text = run.text.replace("5.5", "5")
    print("2. EKF section divider fixed")

    # -----------------------------------------------------------------------
    # 3. Repurpose slide 16 (index 15) as "Cross-City Generalisation: Setup"
    #    instead of deleting it (avoids ZIP duplicate-part issues)
    # -----------------------------------------------------------------------
    def clear_content_shapes(slide, keep_names=None):
        """Remove non-structural text/picture shapes from a slide."""
        if keep_names is None:
            keep_names = set()
        spTree = slide.shapes._spTree
        to_remove = []
        for shape in slide.shapes:
            nm = shape.name
            if nm in keep_names:
                continue
            if nm.startswith("Image") or nm.startswith("Shape"):
                continue
            to_remove.append(shape._element)
        for el in to_remove:
            spTree.remove(el)

    s16 = prs.slides[15]  # blank cross-city slide
    clear_content_shapes(s16)
    add_frame(s16, "CROSS-CITY GENERALISATION — THE CHALLENGE")
    add_box(s16,
            "Can a model trained in Hangzhou and Hong Kong predict GNSS quality in Tokyo — "
            "a city it has never seen?",
            left=0.67, top=1.10, width=12.0, height=0.50,
            size=16, color=DARK, bold=False, align=PP_ALIGN.CENTER)

    challenges = [
        ("Different skyline geometry",
         "Tokyo Shinjuku has a distinct building density and canyon profile vs Hangzhou ring-roads."),
        ("Different satellite visibility",
         "Latitude difference changes GNSS constellation elevation angles and multipath patterns."),
        ("Different receiver hardware",
         "Trimble + u-blox F9P in Tokyo vs Septentrio survey-grade in Beihang — different noise floors."),
        ("Zero training data from Tokyo",
         "Tokyo is completely excluded from training. Not a single Tokyo window was used. "
         "This is a true zero-shot evaluation."),
    ]
    top = 1.80
    for title, body in challenges:
        add_filled_box(s16, 0.67, top, 0.08, 0.80, BLUE_DARK)
        add_box(s16, title, left=0.87, top=top, width=11.80, height=0.32,
                size=13, color=BLUE_DARK, bold=True)
        add_box(s16, body, left=0.87, top=top + 0.33, width=11.80, height=0.38,
                size=11, color=GREY)
        top += 1.05

    add_filled_box(s16, 0.67, 6.40, 12.0, 0.50, NAVY)
    add_box(s16, "Next slide: the actual numbers →  Does SENTINEL hold up?",
            left=0.77, top=6.40, width=11.80, height=0.50,
            size=14, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    print("3. Slide 16 repurposed as Cross-City context")

    # -----------------------------------------------------------------------
    # 4. Rebuild slides 23 and 24 (indices 22 and 23 — no deletion, no shift)
    # -----------------------------------------------------------------------

    # Wipe slide 23 (index 22) → "Why EKF?"
    s23 = prs.slides[22]
    clear_content_shapes(s23)
    for shape in s23.shapes:
        if shape.name == "TextBox 2" and shape.has_text_frame:
            try:
                shape.text_frame.paragraphs[0].runs[0].text = "WHY THE ADAPTIVE EKF?"
            except (IndexError, AttributeError):
                pass
            break

    # Add content to slide 21
    # Left column header
    add_filled_box(s23, 0.67, 1.55, 5.90, 0.42, RGBColor(0xC6, 0x28, 0x28))
    add_box(s23, "Standard Kalman Filter",
            left=0.67, top=1.55, width=5.90, height=0.42,
            size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_filled_box(s23, 6.77, 1.55, 5.90, 0.42, GREEN)
    add_box(s23, "Our Adaptive EKF",
            left=6.77, top=1.55, width=5.90, height=0.42,
            size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    left_items = [
        ("R = fixed (e.g. 9 m²)", RED),
        ("Always trusts GNSS equally", GREY),
        ("Reactive — waits until GNSS fails", GREY),
        ("Position jumps during blockage", GREY),
        ("No advance warning used", GREY),
    ]
    top = 2.10
    for txt, col in left_items:
        add_box(s23, "●  " + txt, left=0.77, top=top, width=5.60, height=0.40,
                size=13, color=col, bold=(col == RED))
        top += 0.47

    right_items = [
        ("R(t) grows with P̂(DEGRADED)", GREEN),
        ("Pre-emptively distrusts GNSS", DARK),
        ("Shifts to dead-reckoning before failure", DARK),
        ("Smooth handoff during blockage", DARK),
        ("5 s early warning from SENTINEL", DARK),
    ]
    top = 2.10
    for txt, col in right_items:
        add_box(s23, "✔  " + txt, left=6.87, top=top, width=5.70, height=0.40,
                size=13, color=col, bold=(col == GREEN))
        top += 0.47

    div = s23.shapes.add_shape(1, Inches(6.60), Inches(
        1.55), Inches(0.05), Inches(3.00))
    div.fill.solid()
    div.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    div.line.fill.background()

    add_filled_box(s23, 0.67, 6.45, 12.0, 0.58, BLUE_DARK)
    add_box(s23,
            "“Prediction closes the loop: we don’t wait for GNSS to fail — we pre-empt it”",
            left=0.77, top=6.45, width=11.80, height=0.58,
            size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    print("4a. Slide 23 (Why EKF?) rebuilt")

    # Wipe slide 24 (index 23) → "Formula"
    s24 = prs.slides[23]
    clear_content_shapes(s24)
    for shape in s24.shapes:
        if shape.name in ("TextBox 4", "Text 1") and shape.has_text_frame:
            shape.text_frame.paragraphs[0].clear()

    add_frame(s24, "ADAPTIVE MEASUREMENT NOISE — THE MECHANISM")

    add_filled_box(s24, 1.00, 1.05, 11.33, 0.80, RGBColor(0xF0, 0xF4, 0xFF))
    add_box(s24, "R(t)  =  σ²_base  +  (σ²_deg − σ²_base)  ×  P̂_calib(t)",
            left=1.00, top=1.05, width=11.33, height=0.80,
            size=24, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

    add_filled_box(s24, 2.00, 1.95, 9.33, 0.55, RGBColor(0xE8, 0xF0, 0xFE))
    add_box(s24, "P̂_calib(t)  =  clip ( ( P̂(t) − P₅ ) / ( 1 − P₅ ),   0,   1 )",
            left=2.00, top=1.95, width=9.33, height=0.55,
            size=17, color=BLUE_DARK, align=PP_ALIGN.CENTER)

    add_box(s24, "Kalman gain:   Kₜ  =  P⁻ₜ Hᵀ ( H P⁻ₜ Hᵀ + Rₜ )⁻¹",
            left=1.00, top=2.62, width=11.33, height=0.45,
            size=15, color=GREY, align=PP_ALIGN.CENTER)

    annotations = [
        ("σ²_base = 9 m²",
         "Baseline GNSS noise (CLEAN signal) — filter trusts GNSS tightly"),
        ("σ²_deg = 10,000 m²",
         "Noise under full degradation — filter ignores GNSS, dead-reckons"),
        ("P̂_calib(t)",
         "Calibrated DEGRADED probability from SENTINEL (0 = clean, 1 = degraded)"),
        ("P₅ = 0.153",
         "Floor offset — unsupervised calibration removes cross-receiver bias"),
        ("Kₜ shrinks when Rₜ rises",
         "Filter leans on motion model when SENTINEL predicts degradation"),
    ]
    top = 3.18
    for term, desc in annotations:
        add_box(s24, term, left=0.77, top=top, width=2.60, height=0.35,
                size=11, color=BLUE_DARK, bold=True)
        add_box(s24, desc, left=3.50, top=top, width=9.30, height=0.35,
                size=11, color=GREY)
        top += 0.42

    pills = [
        (0.67,  "P̂ = 0",   "R = 9 m²  →  Trust GNSS",       GREEN),
        (4.72,  "P̂ = 0.5", "R ≈ 500 m²  →  Caution",    AMBER),
        (8.77,  "P̂ = 1",   "R = 10,000 m²  →  Dead-reckon", RED),
    ]
    for x, label, desc, col in pills:
        add_filled_box(s24, x, 6.30, 3.55, 0.78, col)
        add_box(s24, label, left=x, top=6.30, width=3.55, height=0.35,
                size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_box(s24, desc, left=x, top=6.62, width=3.55, height=0.35,
                size=10, color=WHITE, align=PP_ALIGN.CENTER)

    print("4b. Slide 24 (Formula) rebuilt")

    # -----------------------------------------------------------------------
    # 5. Trim validation slide (index 24, slide 25 — no deletion shift)
    # -----------------------------------------------------------------------
    rebuild_validation_slide(prs.slides[24])
    print("5. Validation slide trimmed to E3/E4/E6/E7")

    # -----------------------------------------------------------------------
    # 6. Add 4 new slides at end, then reorder sldIdLst in one shot
    #    Starting slides: 32 (indices 0-31)
    #    After adding 4:   36 (indices 0-35), new slides at 32-35
    #    Target positions (after formula at index 23):
    #      24 = EKF Results
    #      25 = EKF Severity
    #      26 = Dashboard Overview
    #      27 = Dashboard Demo
    #      28-35 = old slides 24-31 (shifted +4)
    # -----------------------------------------------------------------------
    build_ekf_results(prs)
    build_ekf_severity(prs)
    build_dashboard_overview(prs)
    build_dashboard_demo(prs)

    n_now = len(prs.slides)  # = 36
    print(f"6. Added 4 new slides (total now {n_now})")

    # Move the 4 new slides (at 32-35) into positions 24-27
    # Each move is independent after accounting for the shifting.
    # Use n_now as the base — after each removal the indices of later
    # slides shift down by 1, but since we capture n_now-4..n_now-1
    # we track the correct current index at each step.
    move_slide(prs, n_now - 4, 24)   # ekf_results:  32 → 24
    # ekf_severity: after prev move still at 32 → 25
    move_slide(prs, n_now - 3, 25)
    move_slide(prs, n_now - 2, 26)   # dash_overview: 33 → 26
    move_slide(prs, n_now - 1, 27)   # dash_demo:     34 → 27
    print("7. EKF Results, Severity, Dashboard slides moved to positions 24-27")

    # -----------------------------------------------------------------------
    # 7. Save
    # -----------------------------------------------------------------------
    prs.save(DST)
    print(f"\nSaved: {DST}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
