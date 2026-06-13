"""
update_slide25_ekf_cards.py
----------------------------
Replace the stale synthetic / semi-synthetic EKF-result cards on slide 25 of
docs/SENTINEL_GNSS_Presentation_updated.pptx with the verified, all-real Tokyo
numbers:

  Card 1: +38.7%  Calibrated SENTINEL          47.4 -> 29.0 m   (blue)
  Card 2: +43.2%  + Online sigma_deg (self-tune) 47.4 -> 26.9 m (green)
  Card 3: +48.8%  Fixed-R EKF (best, clean)     47.4 -> 24.3 m  (green)

Edits the _updated copy in place (the original may be open in PowerPoint).
Run: python -m src.utils.update_slide25_ekf_cards
"""
from __future__ import annotations
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

ROOT = Path(__file__).resolve().parents[2]
DECK = ROOT / "docs" / "SENTINEL_GNSS_Presentation_updated.pptx"

BLUE = RGBColor(0x00, 0x38, 0x93)
GREEN = RGBColor(0x1B, 0x87, 0x3A)
NAVY = RGBColor(0x0A, 0x17, 0x33)

# new content keyed by a substring of the CURRENT text in each box
NUMBERS = {
    "33.8": ("+38.7%", BLUE),
    "82%":  ("+43.2%", GREEN),
    "48.8": ("+48.8%", GREEN),
}
CAPTIONS = {
    "Synthetic":      ("Calibrated SENTINEL", "real Tokyo  47.4 -> 29.0 m"),
    "Semi-synthetic": ("+ Online σ_deg (self-tuning)", "real Tokyo  47.4 -> 26.9 m"),
    "Fully real":     ("Fixed-R EKF (best on clean receiver)", "real Tokyo  47.4 -> 24.3 m"),
}


def _first_run_size(tf):
    for p in tf.paragraphs:
        for r in p.runs:
            if r.font.size:
                return r.font.size
    return Pt(12)


def set_number(tf, text, color):
    p = tf.paragraphs[0]
    if not p.runs:
        return
    p.runs[0].text = text
    p.runs[0].font.color.rgb = color
    # drop any extra runs
    for r in p.runs[1:]:
        r._r.getparent().remove(r._r)


def set_caption(tf, label, value):
    size = _first_run_size(tf)
    tf.clear()
    p0 = tf.paragraphs[0]
    r0 = p0.add_run(); r0.text = label
    r0.font.size = size; r0.font.bold = True
    r0.font.color.rgb = NAVY; r0.font.name = "Calibri"
    p1 = tf.add_paragraph()
    r1 = p1.add_run(); r1.text = value
    r1.font.size = size; r1.font.color.rgb = NAVY; r1.font.name = "Calibri"


def main():
    prs = Presentation(str(DECK))
    slide = prs.slides[24]   # slide 25
    n_num = n_cap = 0
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        txt = sh.text_frame.text
        for key, (new, col) in NUMBERS.items():
            if key in txt and len(txt) < 12:   # the big-number boxes are short
                set_number(sh.text_frame, new, col); n_num += 1; break
        else:
            for key, (lab, val) in CAPTIONS.items():
                if key in txt:
                    set_caption(sh.text_frame, lab, val); n_cap += 1; break
    prs.save(str(DECK))
    print(f"updated {n_num} number cards and {n_cap} caption boxes on slide 25")
    print(f"saved -> {DECK.name}")


if __name__ == "__main__":
    main()
