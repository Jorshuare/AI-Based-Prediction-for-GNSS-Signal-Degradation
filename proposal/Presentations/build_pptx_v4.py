"""
Build SENTINEL_GNSS_Proposal_v4.pptx from python-pptx.
Beihang (BUAA) theme · 16:9 · 1.5 line spacing · body >= 18pt.
Run:  python build_pptx_v4.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Beihang palette ──────────────────────────────────────────────────────────
NAVY = RGBColor(0x00, 0x33, 0x66)
BLUE = RGBColor(0x00, 0x5B, 0xAC)
SKY = RGBColor(0x00, 0xA0, 0xE9)
BG = RGBColor(0xF5, 0xF8, 0xFC)
GREY = RGBColor(0x5A, 0x5A, 0x5A)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CLEAN = RGBColor(0x4C, 0xAF, 0x50)
WARN = RGBColor(0xFF, 0x98, 0x00)
DEG = RGBColor(0xF4, 0x43, 0x36)

HERE = os.path.dirname(os.path.abspath(__file__))
FIG_DIR = os.path.abspath(os.path.join(HERE, "..", "..", "results", "figures"))

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf


def _para(tf, text, size=18, color=DARK, bold=False, bullet=False, space_after=6,
          first=False, align=PP_ALIGN.LEFT):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.line_spacing = 1.5
    p.space_after = Pt(space_after)
    p.alignment = align
    run = p.add_run()
    run.text = ("•  " + text) if bullet else text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return p


def title_bar(slide, text, color=NAVY, size=34):
    # accent bar
    bar = slide.shapes.add_shape(
        1, Inches(0.6), Inches(0.5), Inches(0.18), Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    _, tf = _box(slide, 0.95, 0.45, 11.8, 1.0)
    _para(tf, text, size=size, color=color, bold=True, first=True)


def title_slide():
    s = prs.slides.add_slide(BLANK)
    _bg(s, NAVY)
    band = s.shapes.add_shape(1, 0, Inches(2.2), SW, Inches(2.0))
    band.fill.solid()
    band.fill.fore_color.rgb = BLUE
    band.line.fill.background()
    _, tf = _box(s, 0.8, 2.25, 11.8, 1.0)
    _para(tf, "SENTINEL-GNSS", size=50, color=WHITE, bold=True, first=True)
    _, tf2 = _box(s, 0.8, 3.25, 11.8, 1.2)
    _para(tf2, "Proactive Multi-Horizon Prediction of GNSS Signal Degradation "
               "for Autonomous Vehicle Navigation", size=24, color=RGBColor(0xE6, 0xF2, 0xFB),
          bold=False, first=True)
    _, tf3 = _box(s, 0.8, 4.7, 11.8, 1.5)
    _para(tf3, "Team Pilot Project — Beihang University",
          size=22, color=WHITE, bold=True, first=True)
    _para(tf3, "Progress Presentation · Version 4 · Run 14 results",
          size=18, color=RGBColor(0xCF, 0xE3, 0xF5))
    _para(tf3, "Forecasting GNSS degradation 5, 15 and 30 seconds before it happens.",
          size=16, color=RGBColor(0xCF, 0xE3, 0xF5))


def section_divider(num, text):
    s = prs.slides.add_slide(BLANK)
    _bg(s, NAVY)
    _, tf = _box(s, 0.9, 3.0, 11.5, 1.5)
    _para(tf, f"{num} · {text}", size=46, color=WHITE, bold=True, first=True)
    line = s.shapes.add_shape(1, Inches(0.95), Inches(
        4.25), Inches(3.2), Inches(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = SKY
    line.line.fill.background()


def bullets_slide(title, bullets, sizes=None):
    s = prs.slides.add_slide(BLANK)
    _bg(s, BG)
    title_bar(s, title)
    _, tf = _box(s, 0.95, 1.7, 11.6, 5.4)
    for i, b in enumerate(bullets):
        if isinstance(b, tuple):
            txt, col = b
        else:
            txt, col = b, DARK
        sz = sizes[i] if sizes else 20
        _para(tf, txt, size=sz, color=col, bullet=True,
              first=(i == 0), space_after=10)


def table_slide(title, headers, rows, note=None, col_widths=None, fontsize=15):
    s = prs.slides.add_slide(BLANK)
    _bg(s, BG)
    title_bar(s, title)
    nrows, ncols = len(rows) + 1, len(headers)
    tbl_w = 12.0
    tbl = s.shapes.add_table(nrows, ncols, Inches(0.95), Inches(1.75),
                             Inches(tbl_w), Inches(0.4 * nrows)).table
    if col_widths:
        for j, w in enumerate(col_widths):
            tbl.columns[j].width = Inches(w)
    for j, h in enumerate(headers):
        c = tbl.cell(0, j)
        c.text = h
        c.fill.solid()
        c.fill.fore_color.rgb = BLUE
        p = c.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.runs[0]
        r.font.size = Pt(fontsize)
        r.font.bold = True
        r.font.color.rgb = WHITE
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.text = str(val)
            c.fill.solid()
            c.fill.fore_color.rgb = WHITE if i % 2 else RGBColor(
                0xEC, 0xF3, 0xFA)
            p = c.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            rn = p.runs[0]
            rn.font.size = Pt(fontsize)
            rn.font.color.rgb = DARK
    if note:
        _, tf = _box(s, 0.95, 1.75 + 0.42 * nrows + 0.1, 11.8, 1.2)
        _para(tf, note, size=15, color=GREY, first=True)


def figure_slide(title, fig_name, caption):
    s = prs.slides.add_slide(BLANK)
    _bg(s, BG)
    title_bar(s, title)
    # python-pptx cannot embed PDF — always use the PNG sibling.
    if fig_name.lower().endswith(".pdf"):
        fig_name = fig_name[:-4] + ".png"
    pf = os.path.abspath(os.path.join(
        HERE, "..", "..", "results", "paper_figures", fig_name))
    path = pf if os.path.exists(pf) else os.path.join(FIG_DIR, fig_name)
    if os.path.exists(path):
        s.shapes.add_picture(path, Inches(
            2.2), Inches(1.8), height=Inches(4.6))
        cap = caption
    else:
        ph = s.shapes.add_shape(1, Inches(2.5), Inches(
            2.0), Inches(8.3), Inches(4.0))
        ph.fill.solid()
        ph.fill.fore_color.rgb = RGBColor(0xE3, 0xEC, 0xF5)
        ph.line.color.rgb = BLUE
        _para(ph.text_frame, f"[ INSERT FIGURE ]\n{fig_name}", size=20, color=BLUE,
              bold=True, first=True, align=PP_ALIGN.CENTER)
        cap = caption + f"  (drop in results/figures/{fig_name})"
    _, tf = _box(s, 0.95, 6.5, 11.8, 0.8)
    _para(tf, cap, size=15, color=GREY, first=True)


# ════════════════════════════════════════════════════════════════════════════
# BUILD
# ════════════════════════════════════════════════════════════════════════════
title_slide()

section_divider("1", "The Problem")
bullets_slide("Why GNSS prediction matters for autonomous driving", [
    "GNSS is the primary absolute-positioning sensor for autonomous vehicles.",
    "It fails abruptly: urban canyons, tunnels, foliage, multipath.",
    "A car at 60 km/h travels 17 metres every second.",
    ("Every existing monitor (RTKLIB codes, RAIM, ML classifiers) is REACTIVE — "
     "it reports degradation only after it happens.", DEG),
    ("Our question is harder: will the signal degrade in the next 5 / 15 / 30 seconds?", BLUE),
])
table_slide("What proactive warning buys the vehicle",
            ["Horizon", "Distance @ 60 km/h", "Action the planner can take"],
            [["+5 s", "83 m", "Tighten IMU fusion, slow down"],
                ["+15 s", "250 m", "Pre-engage dead-reckoning, adjust speed"],
                ["+30 s", "500 m", "Re-route to avoid the degradation zone"]],
            note="Prediction converts a sudden failure into a planned hand-off to backup localisation.",
            col_widths=[2.2, 3.2, 6.6], fontsize=17)

section_divider("2", "Data")
table_slide("A multi-city, multi-receiver dataset (149,662 labelled epochs)",
            ["Source", "City", "Receivers", "Role"],
            [["Field Scenarios A–E", "Beihang", "Septentrio", "train / test"],
                ["UrbanNav Medium/Tunnel/Deep/Harsh",
                    "Hong Kong", "9+", "train / val"],
                ["Tokyo Shinjuku", "Tokyo", "Trimble + u-blox", "HELD-OUT city"]],
            note="5 controlled degradation scenarios (A–E). 4 cities, 9+ receiver types, one 3-class schema.",
            col_widths=[5.0, 2.2, 2.8, 2.0], fontsize=16)
bullets_slide("From raw signals to model-ready features", [
    "Raw RINEX + NMEA  →  37 engineered features in 7 groups, per 1-second epoch.",
    "Groups: position, C/N0 signal strength, satellite count, DOP, receiver status, "
    "temporal patterns, atmospheric.",
    "A 30-second sliding window → tensor (30 × 37); labels at t+5 / t+15 / t+30 s.",
    ("CLEAN = healthy fix", CLEAN), ("WARNING = partial degradation", WARN),
    ("DEGRADED = loss of fix / severe", DEG),
])

section_divider("3", "Method")
bullets_slide("The SENTINEL-GNSS architecture", [
    "Transformer encoder (2 layers, 8 heads, d=128) — long-range dependencies in the window.",
    "BiLSTM (2 layers, hidden=256) — directional trajectory toward degradation.",
    "Three output heads (+5s, +15s, +30s) + auxiliary head — one model, one forward pass.",
    "1.46 M parameters; focal loss (gamma=1.0) + class weights [1, 2, 5]; NO SMOTE for the network.",
    ("Honest validation: ablations + permutation test + temporal-feature ablation — "
     "we report negative results too.", BLUE),
])
figure_slide("Architecture", "fig14_architecture.png",
             "Block diagram: input 30×37 → projection → Transformer ×2 → BiLSTM ×2 → 3 heads.")

section_divider("4", "Results (Run 15)")
figure_slide("Multi-horizon performance", "fig01_multihorizon.png",
             "Macro-F1 and MCC at +5/+15/+30 s with 95% bootstrap confidence intervals.")
figure_slide("Per-class metrics at +5 s", "fig02_perclass_f1.png",
             "Per-class F1 across horizons; DEGRADED (safety-critical) shown with recall.")
table_slide("Headline: multi-horizon prediction",
            ["Horizon", "Accuracy", "Macro-F1", "MCC", "95% CI (Macro-F1)"],
            [["+5 s", "0.854", "0.821", "0.773", "[0.800, 0.843]"],
                ["+15 s", "0.789", "0.741", "0.691", "[0.717, 0.764]"],
                ["+30 s", "0.830", "0.783", "0.731", "[0.758, 0.804]"]],
            note="DEGRADED recall = 0.85 at +5 s. DEGRADED F1 improved 0.274 → 0.718 across runs (2.6×).",
            col_widths=[2.0, 2.3, 2.3, 2.0, 3.4], fontsize=17)
table_slide("Per-class performance at +5 s",
            ["Class", "Precision", "Recall", "F1", "Support"],
            [["CLEAN", "0.868", "0.993", "0.927", "731"],
                ["WARNING", "0.947", "0.718", "0.817", "746"],
                ["DEGRADED", "0.623", "0.847", "0.718", "209"]],
            note="High recall on the safety-critical DEGRADED class is the priority for an AV system.",
            col_widths=[2.6, 2.4, 2.2, 2.2, 2.6], fontsize=17)
figure_slide("Ablation — component contribution", "fig04_ablation.png",
             "Full vs LSTM-only vs Transformer-only; full model wins on all metrics.")
table_slide("KEY RESULT — cross-city generalisation (Tokyo, never seen)",
            ["Model", "Beihang/Beijing", "Tokyo", "Tokyo DEGRADED F1"],
            [["DL + XGBoost ENSEMBLE", "0.911", "0.892", "0.896"],
                ["SENTINEL-GNSS (DL)", "0.822", "0.649", "0.753"],
                ["XGBoost", "0.919", "0.821", "0.784"],
                ["RandomForest", "0.926", "0.618", "0.148"]],
            note="In-domain trees win. Cross-city, a DL+XGBoost ENSEMBLE beats every single model "
            "(DEGRADED 0.90); RandomForest collapses (0.15) but XGBoost transfers. Ensemble = the "
            "deployment choice feeding the EKF.",
            col_widths=[4.6, 2.6, 2.0, 2.6], fontsize=15)
bullets_slide("Efficiency, ensemble & calibration", [
    ("Cross-city: DL+XGBoost ensemble is best (Macro-F1 0.892, DEGRADED 0.896).", BLUE),
    "Inference: 0.045 ms/sample on GPU — ~9.5x faster than three separate tree models.",
    "One 17.8 MB checkpoint serves all three horizons; real-time at 10 Hz.",
    "Calibration: temperature scaling cuts ECE 40% (0.114 -> 0.069); not yet <0.05 (future work).",
    ("Persistence baseline scores 0.91 at +5s -> the model's value is transitions + longer horizons.", BLUE),
])
figure_slide("Ensemble wins cross-city", "fig16_ensemble.png",
             "In-domain vs cross-city (Tokyo): DL+XGBoost soft-vote beats every single model.")
figure_slide("Cross-city generalization", "fig05_crosscity_degraded.png",
             "RandomForest collapses cross-city (DEGRADED 0.15); XGBoost transfers (0.78); DL (0.75).")
figure_slide("Real-data EKF case study", "fig18_ekf_realdata.png",
             "Raw GNSS trajectory vs adaptive EKF on real Beihang field NMEA; P(DEGRADED) predictions.")

section_divider("5", "Novelty & Validation")
bullets_slide("What is genuinely new", [
    ("1. Reactive → proactive: first multi-horizon GNSS degradation PREDICTOR.", BLUE),
    "2. Cross-city: DL+XGBoost ensemble best (DEGRADED 0.90); RandomForest collapses (0.15), XGBoost transfers.",
    "3. Unified multi-horizon model — one pass, three horizons, ~9.5x faster.",
    "4. Multi-city, multi-receiver open benchmark (149,662 epochs, reproducible).",
    "5. Hardware-aware design — explicit receiver-tier conditioning across 9+ receivers.",
])
bullets_slide("How we defend it to reviewers (honest narrative)", [
    "\"Trees beat you in-domain\" -> true; cross-city a DL+XGBoost ensemble wins (0.892), one model option ~9.5x faster.",
    "\"Does the Transformer use time?\" → order adds ~3% (permutation test); the win is "
    "representation transfer, not order modelling.",
    "\"Small DEGRADED test set\" → bootstrap 95% CIs on every per-class metric.",
    "\"Data leakage?\" → session-level split; scaler/SMOTE on train only; Tokyo fully held out.",
    ("We lead with generalisation + efficiency, not an in-domain leaderboard score.", BLUE),
])

section_divider("6", "Roadmap & Deliverables")
bullets_slide("Publication plan — 2 papers + 1 conference", [
    ("Paper A (GPS Solutions, Q1): method + multi-horizon + cross-receiver + cross-city + EKF.", BLUE),
    "Paper B (Scientific Data / Data in Brief): the benchmark / dataset descriptor.",
    "Conference (ION GNSS+ 2026): cross-city result as a short paper, later extended.",
    "Two substantial papers avoid salami-slicing and carry more impact than four thin ones.",
])
section_divider("6", "Future Work & Deployment")
bullets_slide("Real-data EKF: from simulation to deployed navigation", [
    "Current: adaptive EKF tested on controlled blockage simulation (33.8% RMSE gain).",
    "Next: integrate UrbanNav Tokyo ground-truth (reference.csv, cm-level SPAN-INS).",
    "Compute rover single-point positions (RTKLIB SPP from rover.obs), align with reference.",
    "Real-data case study: demonstrate 5–15% RMSE gain on actual urban degradation events.",
    ("Phase 2: per-receiver EKF tuning; multi-constellation (GPS+BeiDou+Galileo).", BLUE),
])
bullets_slide("Ensemble model deployment & inference pipeline", [
    "Save trained XGBoost model (joblib) during training for reproducibility.",
    "inference.py supports --ensemble flag: load DL + XGB, soft-vote probabilities.",
    "End-to-end: NMEA stream → features → P(DEGRADED) +5/15/30s → EKF trajectory.",
    "Real-time capable: 0.045 ms/sample on GPU; 10 Hz stream at <1% CPU.",
    ("Phase 2b: raw per-satellite C/N₀ streams; +60s horizon retraining.", BLUE),
])
bullets_slide("Web dashboard (Next.js + FastAPI)", [
    "Real-time monitor: C/N₀, DOP, sat count + live 3-horizon prediction bars.",
    "Route map (Mapbox): colour-coded predicted signal quality along the planned path.",
    "Prediction timeline with lead-time annotations; attention heatmap for interpretability.",
    "DL-vs-baseline comparison; dataset explorer (149k epochs); metrics/performance board.",
    ("Phase 2: integrate actual vehicle IMU; Kalman-filter fusion visualization.", BLUE),
])
bullets_slide("Publication & reproducibility", [
    ("Paper A (GPS Solutions, Q1): method + multi-horizon + cross-receiver + cross-city.", BLUE),
    "Paper B (Journal of Navigation): model comparison → ensemble selection → EKF integration.",
    "Conference (ION GNSS+ 2026): cross-city generalization result as a systems paper.",
    "Data + code release: GitHub + Zenodo (reproducible; citable DOI).",
    "All results / figures / models in results/ (zipped & versioned).",
])

# Closing
s = prs.slides.add_slide(BLANK)
_bg(s, NAVY)
_, tf = _box(s, 0.9, 2.6, 11.6, 2.0)
_para(tf, "Thank you", size=46, color=WHITE, bold=True, first=True)
_para(tf, "SENTINEL-GNSS — predicting GNSS degradation before it happens", size=22,
      color=RGBColor(0xE6, 0xF2, 0xFB))
_para(tf, "Team Pilot Project · Beihang University", size=18, color=WHITE)
_para(tf, "github.com/Jorshuare/AI-Based-Prediction-for-GNSS-Signal-Degradation",
      size=15, color=RGBColor(0xCF, 0xE3, 0xF5))

out = os.path.join(HERE, "SENTINEL_GNSS_Proposal_v4.pptx")
prs.save(out)
print(f"Saved {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
