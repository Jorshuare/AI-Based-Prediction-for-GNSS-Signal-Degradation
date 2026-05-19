"""
Build GNSS_Proposal_Presentation_v3.pptx
=========================================
Adds 5 new "AI Model & Training" slides after the V2 progress section,
reflecting:
  • SENTINEL-GNSS Transformer-LSTM architecture
  • Dataset expansion to 97,393 rows (10 sources, 3 cities)
  • Full training pipeline (feature_prep → SMOTE → train → evaluate)
  • Smoke-test result: val macro-F1 = 0.9375 at epoch 3
  • Colab T4 GPU training setup

Run:
    python proposal/build_pptx_v3.py
"""

import copy
import re
from pathlib import Path
from lxml import etree

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SRC = Path("proposal/Presentations/GNSS_Proposal_Presentation_v2.pptx")
DST = Path("proposal/Presentations/GNSS_Proposal_Presentation_v3.pptx")

# ── Brand colours (match existing deck) ──────────────────────────────────────
C_DARK = RGBColor(0x0D, 0x1B, 0x2A)
C_ACCENT = RGBColor(0x00, 0xA8, 0xFF)
C_GREEN = RGBColor(0x00, 0xC8, 0x6E)
C_AMBER = RGBColor(0xFF, 0xA5, 0x00)
C_RED = RGBColor(0xD5, 0x5E, 0x00)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT = RGBColor(0xE8, 0xF4, 0xFF)
C_GREY = RGBColor(0x55, 0x65, 0x75)
C_PURPLE = RGBColor(0x7B, 0x2F, 0xBE)

W = Inches(10.0)
H = Inches(5.625)


# ═══════════════════════════════════════════════════════════════════════════════
#  HELPERS  (same API as v2)
# ═══════════════════════════════════════════════════════════════════════════════

def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    fill = shape.fill
    if fill_rgb:
        fill.solid()
        fill.fore_color.rgb = fill_rgb
    else:
        fill.background()
    line = shape.line
    if line_rgb:
        line.color.rgb = line_rgb
        line.width = line_w
    else:
        line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text_box(slide, text, l, t, w, h,
                 size=Pt(11), bold=False, color=C_WHITE,
                 align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_para(tf, text, size=Pt(10), bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, space_before=Pt(0), italic=False):
    from pptx.oxml.ns import qn
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def slide_title_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, W, Inches(0.75), fill_rgb=C_DARK)
    add_text_box(slide, title,
                 Inches(0.3), Inches(0.1), Inches(8), Inches(0.55),
                 size=Pt(20), bold=True, color=C_WHITE)
    add_rect(slide, 0, Inches(0.75), W, Emu(18000), fill_rgb=C_ACCENT)
    if subtitle:
        add_text_box(slide, subtitle,
                     Inches(0.3), Inches(0.82), Inches(9.4), Inches(0.35),
                     size=Pt(10.5), bold=False, color=C_GREY, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
#  V3 SLIDE 1 — Divider
# ═══════════════════════════════════════════════════════════════════════════════

def make_v3_divider_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    sp = slide.shapes._spTree
    for el in list(sp):
        sp.remove(el)

    add_rect(slide, 0, 0, W, H, fill_rgb=C_DARK)
    add_rect(slide, 0, Inches(2.3), W, Emu(36000), fill_rgb=C_ACCENT)

    # V3 badge
    add_rect(slide, Inches(0.4), Inches(0.35), Inches(1.4), Inches(0.55),
             fill_rgb=C_ACCENT)
    add_text_box(slide, "VERSION 3",
                 Inches(0.42), Inches(0.38), Inches(1.36), Inches(0.5),
                 size=Pt(11), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_text_box(slide, "AI Model Development Update",
                 Inches(0.4), Inches(1.0), Inches(9.2), Inches(0.9),
                 size=Pt(36), bold=True, color=C_WHITE)
    add_text_box(slide, "May 2026  ·  From data collection to trained model",
                 Inches(0.4), Inches(2.0), Inches(9.2), Inches(0.5),
                 size=Pt(16), bold=False, color=C_ACCENT)

    bullets = [
        "✔  97,393 labelled epochs — 10 sources, 3 cities (Beijing · Hong Kong · Tokyo)",
        "✔  SENTINEL-GNSS Transformer-LSTM architecture finalised — 359 K parameters",
        "✔  Full pipeline: feature_prep → SMOTE → sliding windows → train → evaluate",
        "✔  Smoke test passed — val macro-F1 = 0.9375 at epoch 3 (5-epoch debug run)",
        "⬡  Full 150-epoch training on Colab T4 GPU — in progress",
    ]
    y = Inches(2.75)
    for b in bullets:
        col = C_GREEN if b.startswith("✔") else C_AMBER
        add_text_box(slide, b,
                     Inches(0.5), y, Inches(9.0), Inches(0.36),
                     size=Pt(11.5), bold=False, color=col)
        y += Inches(0.4)
    return slide


# ═══════════════════════════════════════════════════════════════════════════════
#  V3 SLIDE 2 — SENTINEL-GNSS Architecture
# ═══════════════════════════════════════════════════════════════════════════════

def make_architecture_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    sp = slide.shapes._spTree
    for el in list(sp):
        sp.remove(el)

    add_rect(slide, 0, 0, W, H, fill_rgb=C_WHITE)
    slide_title_bar(slide,
                    "SENTINEL-GNSS — Transformer-LSTM Architecture",
                    "Hybrid model: spatial attention (Transformer) + temporal memory (LSTM) → 3 prediction horizons")

    # ── Left column: input & transformer ────────────────────────────────────
    # Input box
    add_rect(slide, Inches(0.25), Inches(1.05), Inches(2.3), Inches(0.72),
             fill_rgb=C_DARK, line_rgb=C_ACCENT, line_w=Pt(1.5))
    add_text_box(slide, "Input Window",
                 Inches(0.35), Inches(1.07), Inches(2.1), Inches(0.28),
                 size=Pt(10), bold=True, color=C_ACCENT)
    add_text_box(slide, "(B,  T=30,  F=34)\n30 time-steps × 34 features",
                 Inches(0.35), Inches(1.33), Inches(2.1), Inches(0.42),
                 size=Pt(9), bold=False, color=C_WHITE)

    # Arrow down
    add_rect(slide, Inches(1.28), Inches(1.77), Emu(18000), Inches(0.18),
             fill_rgb=C_ACCENT)

    # Linear projection box
    add_rect(slide, Inches(0.25), Inches(1.95), Inches(2.3), Inches(0.55),
             fill_rgb=C_LIGHT, line_rgb=C_ACCENT, line_w=Pt(1))
    add_text_box(slide, "Linear Projection",
                 Inches(0.35), Inches(1.97), Inches(2.1), Inches(0.25),
                 size=Pt(9.5), bold=True, color=C_DARK)
    add_text_box(slide, "34 → d_model = 64  + Positional Encoding",
                 Inches(0.35), Inches(2.2), Inches(2.1), Inches(0.28),
                 size=Pt(8.5), bold=False, color=C_GREY)

    # Arrow
    add_rect(slide, Inches(1.28), Inches(2.5), Emu(18000), Inches(0.18),
             fill_rgb=C_ACCENT)

    # Transformer Encoder box
    add_rect(slide, Inches(0.25), Inches(2.68), Inches(2.3), Inches(1.05),
             fill_rgb=RGBColor(0x0D, 0x2E, 0x4A), line_rgb=C_ACCENT, line_w=Pt(1.5))
    add_text_box(slide, "Transformer Encoder  (×2 layers)",
                 Inches(0.35), Inches(2.7), Inches(2.1), Inches(0.28),
                 size=Pt(9.5), bold=True, color=C_ACCENT)
    enc_items = [
        "4 attention heads",
        "d_ff = 256  |  GELU",
        "Pre-LN  |  dropout 0.1",
    ]
    for i, item in enumerate(enc_items):
        add_text_box(slide, f"• {item}",
                     Inches(0.42), Inches(2.98) + Inches(0.24)*i,
                     Inches(2.0), Inches(0.25),
                     size=Pt(8.5), bold=False, color=C_WHITE)

    # Arrow
    add_rect(slide, Inches(1.28), Inches(3.73), Emu(18000), Inches(0.18),
             fill_rgb=C_ACCENT)

    # LSTM box
    add_rect(slide, Inches(0.25), Inches(3.91), Inches(2.3), Inches(0.72),
             fill_rgb=C_PURPLE, line_rgb=C_PURPLE, line_w=Pt(1))
    add_text_box(slide, "LSTM  (×2 layers)",
                 Inches(0.35), Inches(3.93), Inches(2.1), Inches(0.28),
                 size=Pt(9.5), bold=True, color=C_WHITE)
    add_text_box(slide, "hidden = 128  |  dropout 0.1",
                 Inches(0.35), Inches(4.19), Inches(2.1), Inches(0.28),
                 size=Pt(8.5), bold=False, color=C_LIGHT)

    # Arrow
    add_rect(slide, Inches(1.28), Inches(4.63), Emu(18000), Inches(0.18),
             fill_rgb=C_ACCENT)

    # ── Right column: 3 output heads ────────────────────────────────────────
    heads = [
        ("5 s ahead",  Inches(2.85), C_GREEN),
        ("15 s ahead", Inches(5.05), C_AMBER),
        ("30 s ahead", Inches(7.25), C_RED),
    ]
    # Horizontal bar
    add_rect(slide, Inches(2.75), Inches(4.8), Inches(7.0), Emu(9000),
             fill_rgb=C_ACCENT)
    for label, x, col in heads:
        add_rect(slide, x, Inches(4.81), Inches(1.9), Inches(0.65),
                 fill_rgb=col)
        add_text_box(slide, f"Head: {label}",
                     x+Inches(0.08), Inches(4.83), Inches(1.74), Inches(0.28),
                     size=Pt(9.5), bold=True, color=C_DARK)
        add_text_box(slide, "128→64→3  (CLEAN/WARN/DEGR)",
                     x+Inches(0.08), Inches(5.1), Inches(1.74), Inches(0.3),
                     size=Pt(8), bold=False, color=C_DARK)

    # ── Right column: top — model stats ─────────────────────────────────────
    add_rect(slide, Inches(2.8), Inches(1.05), Inches(6.95), Inches(1.3),
             fill_rgb=C_LIGHT, line_rgb=C_ACCENT, line_w=Pt(1))
    add_text_box(slide, "Model Summary",
                 Inches(2.9), Inches(1.07), Inches(6.7), Inches(0.3),
                 size=Pt(11), bold=True, color=C_DARK)
    stats = [
        ("Parameters",        "~359,241  (all trainable)"),
        ("Input shape",       "(B, 30, 34)  — 30×1 Hz windows, 34 features"),
        ("Loss function",     "Focal Loss  γ=2.0  |  class weights [1, 3, 5]"),
        ("Optimiser",         "AdamW  lr=1e-3  |  weight_decay=1e-4"),
        ("LR schedule",       "5-epoch linear warmup + cosine decay"),
        ("Mixed precision",   "AMP enabled (GPU)  |  batch_size=256 on T4"),
    ]
    for i, (k, v) in enumerate(stats):
        y = Inches(1.38) + Inches(0.165)*i
        add_text_box(slide, f"{k}:",
                     Inches(2.9), y, Inches(1.7), Inches(0.22),
                     size=Pt(8.5), bold=True, color=C_DARK)
        add_text_box(slide, v,
                     Inches(4.6), y, Inches(5.0), Inches(0.22),
                     size=Pt(8.5), bold=False, color=C_GREY)

    # ── Right column: training config ───────────────────────────────────────
    add_rect(slide, Inches(2.8), Inches(2.45), Inches(6.95), Inches(1.45),
             fill_rgb=RGBColor(0xF5, 0xF5, 0xF5), line_rgb=C_GREY, line_w=Pt(0.5))
    add_text_box(slide, "Training Configuration",
                 Inches(2.9), Inches(2.47), Inches(6.7), Inches(0.3),
                 size=Pt(11), bold=True, color=C_DARK)
    cfg = [
        ("Max epochs",      "150  (early stop patience = 20)"),
        ("Train/Val/Test",  "Session-based split  70 / 15 / 15 %"),
        ("Class imbalance", "SMOTE oversampling on training windows"),
        ("Checkpointing",   "Save every 10 epochs + best val macro-F1"),
        ("Target metric",   "Macro-F1  (equal weight across 3 classes)"),
    ]
    for i, (k, v) in enumerate(cfg):
        y = Inches(2.8) + Inches(0.22)*i
        add_text_box(slide, f"{k}:",
                     Inches(2.9), y, Inches(1.7), Inches(0.22),
                     size=Pt(8.5), bold=True, color=C_DARK)
        add_text_box(slide, v,
                     Inches(4.6), y, Inches(5.0), Inches(0.22),
                     size=Pt(8.5), bold=False, color=C_GREY)

    return slide


# ═══════════════════════════════════════════════════════════════════════════════
#  V3 SLIDE 3 — Dataset Expansion
# ═══════════════════════════════════════════════════════════════════════════════

def make_dataset_expansion_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    sp = slide.shapes._spTree
    for el in list(sp):
        sp.remove(el)

    add_rect(slide, 0, 0, W, H, fill_rgb=C_WHITE)
    slide_title_bar(slide,
                    "Dataset Expansion — 97,393 Labelled Epochs",
                    "3 quality classes: CLEAN / WARNING / DEGRADED  ·  10 sources  ·  3 cities  ·  4 GNSS constellations")

    # ── Class distribution bar chart (text-based) ────────────────────────────
    add_rect(slide, Inches(0.25), Inches(1.0), Inches(4.1), Inches(1.5),
             fill_rgb=C_LIGHT, line_rgb=C_ACCENT, line_w=Pt(1))
    add_text_box(slide, "Label Distribution",
                 Inches(0.35), Inches(1.02), Inches(3.9), Inches(0.28),
                 size=Pt(10), bold=True, color=C_DARK)
    classes = [
        ("CLEAN",    66.8, C_GREEN,  RGBColor(0x00, 0x72, 0xB2)),
        ("WARNING",  20.3, C_AMBER,  RGBColor(0xE6, 0x9F, 0x00)),
        ("DEGRADED", 12.9, C_RED,    RGBColor(0xD5, 0x5E, 0x00)),
    ]
    for i, (name, pct, col, _) in enumerate(classes):
        y = Inches(1.35) + Inches(0.35)*i
        bar_w = Inches(3.5 * pct / 100)
        add_rect(slide, Inches(0.35), y, bar_w, Inches(0.26), fill_rgb=col)
        add_text_box(slide, f"{name}  {pct:.1f}%",
                     Inches(0.38), y + Emu(30000), Inches(3.5), Inches(0.22),
                     size=Pt(9), bold=False, color=C_DARK)

    # ── Dataset source table ────────────────────────────────────────────────
    add_rect(slide, Inches(0.25), Inches(2.65), Inches(9.5), Inches(0.35),
             fill_rgb=C_DARK)
    headers = ["Source", "City / Region", "Receiver", "Rows", "Scenarios"]
    col_x = [0.25, 2.45, 5.0, 6.7, 7.8]
    col_w = [2.2,  2.55, 1.7, 1.1, 1.95]
    for hdr, x, w in zip(headers, col_x, col_w):
        add_text_box(slide, hdr,
                     Inches(x+0.05), Inches(2.67), Inches(w), Inches(0.28),
                     size=Pt(8.5), bold=True, color=C_WHITE)

    rows = [
        ("Septentrio scenarios (A–E)", "Beijing, China",
         "MOSAIC-X5C",   "~6 000",  "A B C D E"),
        ("Supervisor Vehicle (×4)",    "Beijing, China",
         "u-blox F9P",   "~5 000",  "B C D E"),
        ("Supervisor Drone",           "Beijing, China",
         "u-blox F9P",   "~2 000",  "D E"),
        ("NCLT (2 dates)",             "Ann Arbor, USA",
         "Velodyne+IMU", "~14 400", "B D"),
        ("Oxford RobotCar (×4)",       "Oxford, UK",
         "NovAtel",      "~10 800", "B C D"),
        ("UrbanNav HK",                "Hong Kong, China",
         "u-blox ZED",   "~8 200",  "A B C D"),
        ("UrbanNav Tokyo Odaiba",      "Tokyo, Japan",
         "Trimble/u-blox", "~5 500", "C D"),
        ("UrbanNav Tokyo Shinjuku",    "Tokyo, Japan",
         "Trimble/u-blox", "~5 500", "B C"),
        ("Phone A (GNSS Logger)",      "Beijing, China",
         "Smartphone",   "~4 000",  "A B C D E"),
        ("Phone B (GPX/GNSS Logger)",  "Beijing, China",
         "Smartphone",   "~4 000",  "A B C D E"),
    ]
    for r_i, (src, city, rx, rows_n, scen) in enumerate(rows):
        bg = C_LIGHT if r_i % 2 == 0 else C_WHITE
        row_y = Inches(3.0) + Inches(0.26)*r_i
        add_rect(slide, Inches(0.25), row_y, Inches(9.5), Inches(0.25),
                 fill_rgb=bg, line_rgb=RGBColor(0xDD, 0xDD, 0xDD), line_w=Pt(0.3))
        vals = [src, city, rx, rows_n, scen]
        for val, x, w in zip(vals, col_x, col_w):
            add_text_box(slide, val,
                         Inches(x+0.05), row_y +
                         Emu(8000), Inches(w), Inches(0.22),
                         size=Pt(8), bold=False, color=C_DARK)

    add_text_box(slide,
                 "Total: 97,393 rows  ·  41 features  ·  Labelled CSV tracked in GitHub",
                 Inches(0.25), Inches(5.27), Inches(9.5), Inches(0.25),
                 size=Pt(9.5), bold=True, color=C_ACCENT)
    return slide


# ═══════════════════════════════════════════════════════════════════════════════
#  V3 SLIDE 4 — Training Pipeline
# ═══════════════════════════════════════════════════════════════════════════════

def make_pipeline_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    sp = slide.shapes._spTree
    for el in list(sp):
        sp.remove(el)

    add_rect(slide, 0, 0, W, H, fill_rgb=C_WHITE)
    slide_title_bar(slide,
                    "SENTINEL-GNSS — Full Training Pipeline",
                    "3-command workflow: feature_prep → train → evaluate  ·  resumable, reproducible, debug-mode available")

    # Pipeline steps as connected boxes
    steps = [
        ("1",  "feature_prep",
         "src.models.feature_prep",
         [
             "Loads labelled CSV (97 k rows)",
             "Imputes & scales 34 features",
             "Sliding windows T=30 @ 1 Hz",
             "SMOTE on training split",
             "Saves train/val/test .npz",
         ],
         C_ACCENT),
        ("2",  "train",
         "src.models.train",
         [
             "Builds SentinelGNSS model",
             "AdamW + cosine LR schedule",
             "Focal Loss (γ=2, cls wts)",
             "Saves checkpoint every 10 ep",
             "--resume from last checkpoint",
         ],
         C_GREEN),
        ("3",  "evaluate",
         "src.models.evaluate",
         [
             "Loads best checkpoint",
             "14 analysis figures",
             "Confusion matrix × 3 horizons",
             "Precision-recall & ROC",
             "Per-source performance",
         ],
         C_PURPLE),
    ]

    box_w = Inches(2.8)
    box_h = Inches(3.7)
    arrow_w = Inches(0.45)
    total_w = 3 * float(box_w) + 2 * float(arrow_w)
    start_x = (float(W) - total_w) / 2

    for i, (num, name, module, items, col) in enumerate(steps):
        x = start_x + i * (float(box_w) + float(arrow_w))

        # Main box
        add_rect(slide, x, Inches(1.0), box_w, box_h,
                 fill_rgb=C_DARK, line_rgb=col, line_w=Pt(2))

        # Step number badge
        add_rect(slide, x + Inches(0.08), Inches(1.08), Inches(0.36), Inches(0.36),
                 fill_rgb=col)
        add_text_box(slide, num,
                     x + Inches(0.09), Inches(1.09), Inches(0.34), Inches(0.33),
                     size=Pt(11), bold=True, color=C_DARK, align=PP_ALIGN.CENTER)

        # Module name
        add_text_box(slide, f"python -m {module}",
                     x + Inches(0.08), Inches(1.52), box_w -
                     Inches(0.16), Inches(0.26),
                     size=Pt(7.5), bold=False, color=col, italic=True)

        # Title
        add_text_box(slide, name,
                     x + Inches(0.08), Inches(1.08), box_w -
                     Inches(0.16), Inches(0.4),
                     size=Pt(13), bold=True, color=C_WHITE, align=PP_ALIGN.RIGHT)

        # Bullet items
        for j, item in enumerate(items):
            add_text_box(slide, f"• {item}",
                         x + Inches(0.14), Inches(1.85) + Inches(0.47)*j,
                         box_w - Inches(0.22), Inches(0.42),
                         size=Pt(9), bold=False, color=C_LIGHT)

        # Arrow between boxes
        if i < 2:
            arrow_x = x + float(box_w) + Inches(0.05)
            add_rect(slide, arrow_x, Inches(2.65), arrow_w - Inches(0.1), Inches(0.12),
                     fill_rgb=C_ACCENT)
            add_text_box(slide, "▶",
                         arrow_x + arrow_w/2 -
                         Inches(0.15), Inches(2.55), Inches(0.3), Inches(0.3),
                         size=Pt(14), bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

    # Footer
    add_rect(slide, Inches(0.25), Inches(4.85), Inches(9.5), Inches(0.62),
             fill_rgb=C_LIGHT, line_rgb=C_ACCENT, line_w=Pt(1))
    add_text_box(slide, "Debug mode:",
                 Inches(0.35), Inches(4.9), Inches(1.1), Inches(0.22),
                 size=Pt(9), bold=True, color=C_DARK)
    add_text_box(slide,
                 "python -m src.models.feature_prep --debug   "
                 "python -m src.models.train --debug   "
                 "python -m src.models.evaluate --debug",
                 Inches(1.45), Inches(4.9), Inches(8.2), Inches(0.22),
                 size=Pt(9), bold=False, color=C_GREY, italic=True)
    add_text_box(slide,
                 "500 rows/source  ·  5 epochs  ·  separate debug windows & checkpoints  ·  never overwrites real data",
                 Inches(0.35), Inches(5.12), Inches(9.3), Inches(0.22),
                 size=Pt(8.5), bold=False, color=C_GREY, italic=True)
    return slide


# ═══════════════════════════════════════════════════════════════════════════════
#  V3 SLIDE 5 — Smoke-Test Results + Colab Setup
# ═══════════════════════════════════════════════════════════════════════════════

def make_results_colab_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    sp = slide.shapes._spTree
    for el in list(sp):
        sp.remove(el)

    add_rect(slide, 0, 0, W, H, fill_rgb=C_WHITE)
    slide_title_bar(slide,
                    "Preliminary Results & Colab Training Setup",
                    "Smoke test (5 epochs, debug dataset)  ·  Full 150-epoch run on Colab T4 — in progress")

    # ── Left: smoke test results ─────────────────────────────────────────────
    add_rect(slide, Inches(0.25), Inches(1.0), Inches(4.55), Inches(4.45),
             fill_rgb=C_LIGHT, line_rgb=C_ACCENT, line_w=Pt(1.5))
    add_text_box(slide, "Smoke Test Results  (debug mode)",
                 Inches(0.35), Inches(1.02), Inches(4.35), Inches(0.3),
                 size=Pt(11), bold=True, color=C_DARK)
    add_text_box(slide, "5 epochs  ·  500 rows/source  ·  CPU (Quadro M1200)",
                 Inches(0.35), Inches(1.32), Inches(4.35), Inches(0.22),
                 size=Pt(8.5), italic=True, color=C_GREY)

    # Headline metric
    add_rect(slide, Inches(0.35), Inches(1.62), Inches(4.35), Inches(0.75),
             fill_rgb=C_GREEN)
    add_text_box(slide, "Best val macro-F1",
                 Inches(0.45), Inches(1.65), Inches(2.2), Inches(0.28),
                 size=Pt(10), bold=True, color=C_DARK)
    add_text_box(slide, "0.9375",
                 Inches(2.65), Inches(1.62), Inches(1.9), Inches(0.7),
                 size=Pt(28), bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    add_text_box(slide, "at epoch 3  (early stop patience = 20)",
                 Inches(0.45), Inches(1.9), Inches(2.2), Inches(0.3),
                 size=Pt(8.5), color=C_DARK)

    detail = [
        ("Training time (5 ep)",  "~45 s on CPU"),
        ("Checkpoints saved",     "epoch_3.pt  +  best.pt"),
        ("Figures generated",     "14 evaluation charts"),
        ("Classes resolved",      "CLEAN · WARNING · DEGRADED"),
        ("Pipeline status",       "All 3 commands passed ✔"),
    ]
    for i, (k, v) in enumerate(detail):
        y = Inches(2.47) + Inches(0.42)*i
        add_rect(slide, Inches(0.35), y, Inches(4.35), Inches(0.38),
                 fill_rgb=(C_WHITE if i %
                           2 == 0 else RGBColor(0xF0, 0xF8, 0xFF)),
                 line_rgb=RGBColor(0xCC, 0xCC, 0xCC), line_w=Pt(0.3))
        add_text_box(slide, k,
                     Inches(0.43), y+Emu(8000), Inches(2.0), Inches(0.3),
                     size=Pt(9), bold=True, color=C_DARK)
        add_text_box(slide, v,
                     Inches(2.43), y+Emu(8000), Inches(2.2), Inches(0.3),
                     size=Pt(9), bold=False, color=C_GREY)

    # ── Right: Colab setup ───────────────────────────────────────────────────
    add_rect(slide, Inches(5.05), Inches(1.0), Inches(4.7), Inches(4.45),
             fill_rgb=C_DARK, line_rgb=C_ACCENT, line_w=Pt(1.5))
    add_text_box(slide, "Google Colab Training (Full Run)",
                 Inches(5.15), Inches(1.02), Inches(4.5), Inches(0.3),
                 size=Pt(11), bold=True, color=C_ACCENT)

    colab_items = [
        ("GPU",          "Tesla T4  —  15.6 GB VRAM",           C_GREEN),
        ("Batch size",   "256  (vs 64 on local CPU)",            C_WHITE),
        ("Speed",        "~20× faster than local",              C_GREEN),
        ("Epochs",       "150  (early stop at plateau)",        C_WHITE),
        ("Data",         "CSV cloned from GitHub — no Drive",   C_WHITE),
        ("Checkpoints",  "→ Google Drive (survives disconnect)", C_AMBER),
        ("Figures",      "→ Google Drive (14 PDF + PNG)",       C_AMBER),
        ("Resume",       "--resume flag for auto-recovery",     C_WHITE),
    ]
    for i, (k, v, col) in enumerate(colab_items):
        y = Inches(1.42) + Inches(0.38)*i
        add_text_box(slide, f"{k}:",
                     Inches(5.15), y, Inches(1.3), Inches(0.3),
                     size=Pt(9), bold=True, color=C_ACCENT)
        add_text_box(slide, v,
                     Inches(6.45), y, Inches(3.2), Inches(0.3),
                     size=Pt(9), bold=False, color=col)

    # Command
    add_rect(slide, Inches(5.15), Inches(4.5), Inches(4.5), Inches(0.38),
             fill_rgb=RGBColor(0x1A, 0x2A, 0x3A))
    add_text_box(slide,
                 "!python -m src.models.train --resume --batch_size 256",
                 Inches(5.22), Inches(4.52), Inches(4.36), Inches(0.28),
                 size=Pt(8), bold=False, color=C_GREEN, italic=True)

    # Footer
    add_text_box(slide,
                 "Full results (confusion matrices, per-source breakdown, 14 figures) available after Colab run completes.",
                 Inches(0.25), Inches(5.27), Inches(9.5), Inches(0.22),
                 size=Pt(8.5), bold=False, color=C_GREY, italic=True)
    return slide


# ═══════════════════════════════════════════════════════════════════════════════
#  MAIN
# ═══════════════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation(str(SRC))

    n_v2 = len(prs.slides)  # V2 slide count (should be 27)

    # Add 5 new V3 slides (appended at end)
    make_v3_divider_slide(prs)        # n_v2 + 0
    make_architecture_slide(prs)      # n_v2 + 1
    make_dataset_expansion_slide(prs)  # n_v2 + 2
    make_pipeline_slide(prs)          # n_v2 + 3
    make_results_colab_slide(prs)     # n_v2 + 4

    new_slides = list(range(n_v2, n_v2 + 5))

    # Reorder: insert after the last V2 progress slide (originally after slide 13,
    # i.e. insert after index 12 in original V2 order, so after V2 progress block)
    # V2 inserted 4 progress slides after original slide 9 (index 8).
    # V2 total = 23 original + 4 new = 27 slides (indices 0-26).
    # V2 progress slides are at indices 9,10,11,12 (after reorder in v2).
    # We want V3 slides to follow immediately: insert after index 12.
    original = list(range(n_v2))
    new_order = original[:13] + new_slides + original[13:]
    # Reorder
    xml_slides = prs.slides._sldIdLst
    sld_ids = list(xml_slides)
    reordered = [sld_ids[i] for i in new_order]
    for el in sld_ids:
        xml_slides.remove(el)
    for el in reordered:
        xml_slides.append(el)

    DST.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(DST))
    print(f"Saved: {DST}")
    print(f"Total slides: {len(prs.slides)}")
    for i, s in enumerate(prs.slides):
        for shape in s.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.paragraphs[0].text.strip()
                if t:
                    print(f"  {i+1:2d}. {t[:72]}")
                    break
        else:
            print(f"  {i+1:2d}. (no text)")


if __name__ == "__main__":
    main()
