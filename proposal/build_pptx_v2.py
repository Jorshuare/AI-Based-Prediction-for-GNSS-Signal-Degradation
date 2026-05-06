"""
Build GNSS_Proposal_Presentation_v2.pptx
=========================================
Inserts 4 new "progress update" slides after slide 9 and updates
the Data Sources slide (7) and Timeline slide (22) to reflect
what has been accomplished as of May 2026.

Run:
    python proposal/build_pptx_v2.py
"""

import copy
import re
from pathlib import Path
from lxml import etree

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SRC = Path("proposal/Presentations/GNSS_Proposal_Presentation.pptx")
DST = Path("proposal/Presentations/GNSS_Proposal_Presentation_v2.pptx")

# ── Brand colours (match existing deck) ──────────────────────────────────────
C_DARK = RGBColor(0x0D, 0x1B, 0x2A)   # near-black navy
C_ACCENT = RGBColor(0x00, 0xA8, 0xFF)   # bright blue
C_GREEN = RGBColor(0x00, 0xC8, 0x6E)   # success green
C_AMBER = RGBColor(0xFF, 0xA5, 0x00)   # warning amber
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT = RGBColor(0xE8, 0xF4, 0xFF)   # very light blue bg
C_GREY = RGBColor(0x55, 0x65, 0x75)

W = Inches(10.0)
H = Inches(5.625)


# ═══════════════════════════════════════════════════════════════════════════════
#  HELPERS
# ═══════════════════════════════════════════════════════════════════════════════

def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_w=Pt(0)):
    from pptx.util import Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape = slide.shapes.add_shape(1, l, t, w, h)
    fill = shape.fill
    if fill_rgb:
        fill.solid()
        fill.fore_color.rgb = fill_rgb
    else:
        fill.background()
    line = shape.line
    if line_rgb:
        line.color.rgb = line_rgb
        line.width = line_w
    else:
        line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text_box(slide, text, l, t, w, h,
                 size=Pt(11), bold=False, color=C_WHITE,
                 align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_para(tf, text, size=Pt(10), bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, space_before=Pt(0), italic=False):
    """Add a paragraph to an existing text frame."""
    from pptx.oxml.ns import qn
    from pptx.util import Pt as _Pt
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def slide_title_bar(slide, title, subtitle=None):
    """Standard slide header: dark bar at top with title."""
    # Top bar
    bar = add_rect(slide, 0, 0, W, Inches(0.75), fill_rgb=C_DARK)
    # Title text
    add_text_box(slide, title,
                 Inches(0.3), Inches(0.1), Inches(8), Inches(0.55),
                 size=Pt(20), bold=True, color=C_WHITE)
    # Accent line
    add_rect(slide, 0, Inches(0.75), W, Emu(18000), fill_rgb=C_ACCENT)
    if subtitle:
        add_text_box(slide, subtitle,
                     Inches(0.3), Inches(0.82), Inches(9.4), Inches(0.35),
                     size=Pt(10.5), bold=False, color=C_GREY,
                     italic=True)


def copy_slide(prs_src, slide_src, prs_dst):
    """
    Copy a slide from prs_src into prs_dst (appended at end).
    Returns the new slide in prs_dst.
    """
    template = prs_dst.slides.add_slide(prs_dst.slide_layouts[0])
    # Remove placeholder shapes from template
    sp = template.shapes._spTree
    for el in list(sp):
        sp.remove(el)

    # Copy spTree content
    src_sp = slide_src.shapes._spTree
    for el in src_sp:
        sp.append(copy.deepcopy(el))

    # Copy background
    if slide_src.background.fill.type is not None:
        pass  # background inherits from theme; skip for simplicity

    return template


def reorder_slides(prs, new_order):
    """
    Reorder slides in prs according to new_order (list of 0-based indices
    into the *current* slide list).
    """
    xml_slides = prs.slides._sldIdLst
    sld_ids = list(xml_slides)
    reordered = [sld_ids[i] for i in new_order]
    for el in sld_ids:
        xml_slides.remove(el)
    for el in reordered:
        xml_slides.append(el)


# ═══════════════════════════════════════════════════════════════════════════════
#  NEW SLIDES
# ═══════════════════════════════════════════════════════════════════════════════

def make_divider_slide(prs):
    """Section divider: 'Progress Update — V2 · May 2026'"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    sp = slide.shapes._spTree
    for el in list(sp):
        sp.remove(el)

    # Full background
    add_rect(slide, 0, 0, W, H, fill_rgb=C_DARK)
    # Accent stripe
    add_rect(slide, 0, Inches(2.3), W, Emu(36000), fill_rgb=C_ACCENT)
    # V2 badge
    badge = add_rect(slide, Inches(0.4), Inches(0.35), Inches(1.4), Inches(0.55),
                     fill_rgb=C_ACCENT)
    add_text_box(slide, "VERSION 2",
                 Inches(0.42), Inches(0.38), Inches(1.36), Inches(0.5),
                 size=Pt(11), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    # Main title
    add_text_box(slide, "Progress Update",
                 Inches(0.4), Inches(1.0), Inches(9.2), Inches(0.9),
                 size=Pt(38), bold=True, color=C_WHITE)
    # Date
    add_text_box(slide, "May 2026  ·  What has been accomplished",
                 Inches(0.4), Inches(2.0), Inches(9.2), Inches(0.5),
                 size=Pt(16), bold=False, color=C_ACCENT)
    # Bullet summary
    bullets = [
        "✔  Field data collected — all 5 scenarios, Septentrio MOSAIC-X5C receiver",
        "✔  Public datasets downloaded — NCLT (2 dates), Oxford (4 traversals), UrbanNav HK+Tokyo",
        "✔  Full processing pipeline implemented and validated",
        "✔  Signal quality analysis charts generated for every run",
        "⬡  RTKLIB post-processing → feature extraction → model training  (next steps)",
    ]
    y = Inches(2.75)
    for b in bullets:
        col = C_GREEN if b.startswith("✔") else C_AMBER
        add_text_box(slide, b,
                     Inches(0.5), y, Inches(9.0), Inches(0.36),
                     size=Pt(11.5), bold=False, color=col)
        y += Inches(0.4)
    return slide


def make_field_collection_slide(prs):
    """Slide: Field Data Collection — Complete"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    sp = slide.shapes._spTree
    for el in list(sp):
        sp.remove(el)

    add_rect(slide, 0, 0, W, H, fill_rgb=C_WHITE)
    slide_title_bar(slide,
                    "Field Data Collection — Complete  ✓  (May 5, 2026)",
                    "Septentrio MOSAIC-X5C receiver · Beihang University campus · ~40 km/h max")

    # Receiver info box
    rx_box = add_rect(slide, Inches(0.25), Inches(1.2), Inches(3.4), Inches(1.8),
                      fill_rgb=C_LIGHT, line_rgb=C_ACCENT, line_w=Pt(1))
    add_text_box(slide, "Receiver Specifications",
                 Inches(0.35), Inches(1.25), Inches(3.2), Inches(0.35),
                 size=Pt(10), bold=True, color=C_DARK)
    specs = [
        "Model:      Septentrio MOSAIC-X5C",
        "Serial:     3640133  |  FW 4.12.1",
        "Rate:       1 Hz  |  4 constellations",
        "Constellations: GPS  Galileo  GLONASS  BeiDou",
        "Output:     RINEX 3 (.26O .26N .26G .26P)",
        "Converter:  sbf2rin-15.10.3 (Septentrio RxTools)",
    ]
    for i, s in enumerate(specs):
        add_text_box(slide, s,
                     Inches(0.35), Inches(1.62) + Inches(0.24)*i,
                     Inches(3.2), Inches(0.25),
                     size=Pt(9), color=C_DARK)

    # Scenario table header
    add_rect(slide, Inches(3.85), Inches(1.2), Inches(5.9), Inches(0.35),
             fill_rgb=C_DARK)
    headers = ["Scenario", "Description", "Runs",
               "Duration", "~Epochs", "Mean C/N₀"]
    col_x = [3.85, 4.45, 6.25, 6.85, 7.7, 8.55]
    col_w = [0.6,  1.8,  0.6,  0.85, 0.85, 1.2]
    for hdr, x, w in zip(headers, col_x, col_w):
        add_text_box(slide, hdr,
                     Inches(x+0.04), Inches(1.22), Inches(w), Inches(0.3),
                     size=Pt(8.5), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    rows = [
        ("A", "Instant Blockage",    "Run_1 A2 A3",
         "45–90 s/run", "~180",  "37.2 dB-Hz"),
        ("B", "Urban Canyon",        "Run_1 B2",
         "8–9 min/run", "~1 600", "37.2 dB-Hz"),
        ("C", "Partial Blockage",    "Run_1 C2",
         "10–16 min",   "~2 700", "38.6 dB-Hz"),
        ("D", "Open Sky (baseline)", "Run_1",
         "~10 min",     "~600",  "41.2 dB-Hz ✔"),
        ("E", "Approaching Block.",  "E2 E3",
         "6–7 min/run", "~1 000", "38.6 dB-Hz"),
    ]
    for r_i, (scen, desc, runs, dur, epochs, cnr) in enumerate(rows):
        bg = C_LIGHT if r_i % 2 == 0 else C_WHITE
        add_rect(slide,
                 Inches(3.85), Inches(1.55) + Inches(0.45)*r_i,
                 Inches(5.9), Inches(0.44),
                 fill_rgb=bg, line_rgb=RGBColor(0xCC, 0xCC, 0xCC), line_w=Pt(0.5))
        cells = [scen, desc, runs, dur, epochs, cnr]
        for val, x, w in zip(cells, col_x, col_w):
            col = C_ACCENT if val in ("A", "B", "C", "D", "E") else C_DARK
            bold = val in ("A", "B", "C", "D", "E")
            add_text_box(slide, val,
                         Inches(x+0.04), Inches(1.57) + Inches(0.45)*r_i,
                         Inches(w), Inches(0.38),
                         size=Pt(9), bold=bold, color=col, align=PP_ALIGN.CENTER)

    # Footer note
    add_text_box(slide,
                 "11 RINEX files total  ·  All processed and verified  ·  Charts: results/scenario_analysis/",
                 Inches(0.25), Inches(5.2), Inches(9.5), Inches(0.3),
                 size=Pt(9), bold=False, color=C_GREY, italic=True)
    return slide


def make_signal_quality_slide(prs):
    """Slide: Signal Quality Analysis — Key Findings"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    sp = slide.shapes._spTree
    for el in list(sp):
        sp.remove(el)

    add_rect(slide, 0, 0, W, H, fill_rgb=C_WHITE)
    slide_title_bar(slide,
                    "Signal Quality Analysis — Key Findings",
                    "Derived from RINEX 3 SNR indicators · all 4 constellations · 1 Hz · Septentrio MOSAIC-X5C")

    findings = [
        ("D — Open Sky",
         "41.2 dB-Hz mean  ·  0.5% epochs below 25 dB-Hz  ·  ~39 obs/epoch (most satellites of any scenario)",
         C_GREEN),
        ("C — Partial Blockage (Trees)",
         "38.6 dB-Hz mean  ·  0.6% epochs degraded  ·  Gradual C/N₀ fluctuations consistent with foliage attenuation",
         RGBColor(0x00, 0x8B, 0x4A)),
        ("B — Urban Canyon",
         "37.2 dB-Hz mean  ·  3.0% epochs below threshold  ·  Periodic sharp dips from building reflections visible",
         C_AMBER),
        ("E — Approaching Blockage",
         "38.6 dB-Hz mean  ·  5.9% degraded  ·  Gradual deterioration pattern clearly identifiable for prediction",
         RGBColor(0xFF, 0x70, 0x00)),
        ("A — Instant Blockage",
         "Sharp satellite dropout in first 15 s  ·  Count collapses 24 → 0–2 sats  ·  C/N₀ falls below 25 dB-Hz immediately",
         RGBColor(0xE0, 0x20, 0x20)),
    ]

    for i, (scen, text, col) in enumerate(findings):
        y = Inches(1.1) + Inches(0.84)*i
        add_rect(slide, Inches(0.25), y, Inches(
            0.32), Inches(0.6), fill_rgb=col)
        add_text_box(slide, scen,
                     Inches(0.65), y+Inches(0.03), Inches(3.2), Inches(0.3),
                     size=Pt(11), bold=True, color=C_DARK)
        add_text_box(slide, text,
                     Inches(0.65), y+Inches(0.3), Inches(9.0), Inches(0.35),
                     size=Pt(9.5), bold=False, color=C_GREY)

    add_rect(slide, Inches(0.25), Inches(5.18), Inches(9.5), Emu(6000),
             fill_rgb=C_ACCENT)
    add_text_box(slide,
                 "Key insight: Scenario A shows clear distinct signature; D and A are near-perfectly separable. "
                 "B/C/E have partial overlaps — exactly the challenge the Transformer-LSTM must solve.",
                 Inches(0.35), Inches(5.15), Inches(9.3), Inches(0.38),
                 size=Pt(9.5), bold=False, color=C_DARK, italic=True)
    return slide


def make_dataset_inventory_slide(prs):
    """Slide: Dataset Inventory — All Sources Ready"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    sp = slide.shapes._spTree
    for el in list(sp):
        sp.remove(el)

    add_rect(slide, 0, 0, W, H, fill_rgb=C_WHITE)
    slide_title_bar(slide,
                    "Dataset Inventory — All Sources Ready",
                    "No further downloads required  ·  RTKLIB post-processing is the next action")

    headers = ["Dataset", "Location", "RTKLIB?",
               "Epochs est.", "Scenarios", "Status"]
    col_x = [0.25, 2.5,  5.1,  6.0,   7.1,   8.5]
    col_w = [2.25, 2.6,  0.9,  1.1,   1.4,   1.4]

    add_rect(slide, Inches(0.25), Inches(0.95), Inches(9.5), Inches(0.38),
             fill_rgb=C_DARK)
    for hdr, x, w in zip(headers, col_x, col_w):
        add_text_box(slide, hdr,
                     Inches(x+0.04), Inches(0.97), Inches(w), Inches(0.32),
                     size=Pt(8.5), bold=True, color=C_WHITE)

    datasets = [
        ("Our Collection (Septentrio)",  "raw/scenarios/",
         "YES",  "~6 000",  "A B C D E", "✔ Collected"),
        ("Supervisor Vehicle (exp1–4)",   "raw/supervisor/vehicle/",
         "YES",  "~5 000",  "B C D E",   "✔ On disk"),
        ("Supervisor Drone",             "raw/supervisor/drone/",
         "YES",  "~2 000",  "D E",       "✔ On disk"),
        ("NCLT (2 dates)",               "raw/public/nclt/",
         "NO",   "~14 400", "B D",       "✔ Downloaded"),
        ("Oxford RobotCar (4 traversals)", "raw/public/oxford/",
         "NO",   "~10 800", "B C D",     "✔ Downloaded"),
        ("UrbanNav HK + Tokyo",          "raw/public/urbannav/",
         "YES",  "~5 000",  "A B C D",   "✔ Extracted"),
    ]

    for r_i, (name, loc, rtk, eps, scen, status) in enumerate(datasets):
        bg = C_LIGHT if r_i % 2 == 0 else C_WHITE
        row_y = Inches(1.33) + Inches(0.62)*r_i
        add_rect(slide, Inches(0.25), row_y, Inches(9.5), Inches(0.58),
                 fill_rgb=bg, line_rgb=RGBColor(0xCC, 0xCC, 0xCC), line_w=Pt(0.5))
        vals = [name, loc, rtk, eps, scen, status]
        for val, x, w in zip(vals, col_x, col_w):
            is_status = (val == status)
            col = C_GREEN if "✔" in val else C_DARK
            add_text_box(slide, val,
                         Inches(x+0.04), row_y+Inches(0.04),
                         Inches(w), Inches(0.52),
                         size=Pt(9), bold=(is_status and "✔" in val), color=col)

    add_text_box(slide,
                 f"Total estimated labelled epochs: ~43 000  ·  All 5 scenario classes covered",
                 Inches(0.25), Inches(5.2), Inches(9.5), Inches(0.3),
                 size=Pt(9.5), bold=True, color=C_ACCENT)
    return slide


# ═══════════════════════════════════════════════════════════════════════════════
#  UPDATE EXISTING SLIDES
# ═══════════════════════════════════════════════════════════════════════════════

def update_data_sources_slide(slide):
    """
    Slide 7: Update 'Survey Cart → Data to be acquired'
    to 'Septentrio Receiver → ✔ Data acquired (May 2026)'
    Also update public datasets footnote.
    """
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        for para in tf.paragraphs:
            for run in para.runs:
                if "Data to be acquired" in run.text:
                    run.text = run.text.replace("Data to be acquired",
                                                "✔ Collected (May 2026)")
                if "Survey Cart" in run.text:
                    run.text = run.text.replace("Survey Cart",
                                                "Septentrio MOSAIC-X5C")
                if "Self-collected (5 scenarios)" in run.text:
                    run.text = run.text.replace("Self-collected (5 scenarios)",
                                                "All 5 scenarios collected")
                if "KAIST" in run.text and "NCLT" not in run.text:
                    run.text = run.text.replace(
                        "KAIST",
                        "NCLT ✔  ·  Oxford RobotCar ✔  ·  KAIST (deferred)")


def update_timeline_slide(slide):
    """
    Slide 22: Mark Week 1 tasks as ✅ DONE and add a progress banner.
    """
    week1_tasks = [
        "All software installed",
        "Room 3058",
        "GitHub repo created",
        "ROSBAG extraction",
        "RTKLIB running",
        "Scenarios A & D collected",
        "RTKLIB processed",
    ]
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        for para in tf.paragraphs:
            for run in para.runs:
                for task in week1_tasks:
                    if task in run.text and "✅" not in run.text:
                        run.text = "✅ " + run.text
                if "Scenarios A & D collected" in run.text:
                    run.text = run.text.replace(
                        "Scenarios A & D collected (5 reps each)",
                        "✅ ALL scenarios A–E collected (May 5, 2026)")
                if "Data to be acquired" in run.text:
                    run.text = run.text.replace("Data to be acquired",
                                                "✅ Data collected")


# ═══════════════════════════════════════════════════════════════════════════════
#  MAIN
# ═══════════════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation(str(SRC))

    # ── 1. Update existing slides in-place ──────────────────────────────────
    update_data_sources_slide(prs.slides[6])   # Slide 7
    update_timeline_slide(prs.slides[21])      # Slide 22

    # ── 2. Add 4 new slides (appended at end, indices 23–26) ────────────────
    make_divider_slide(prs)           # will be index 23
    make_field_collection_slide(prs)  # will be index 24
    make_signal_quality_slide(prs)    # will be index 25
    make_dataset_inventory_slide(prs)  # will be index 26

    # ── 3. Reorder: insert new slides after slide 9 (index 8) ───────────────
    # Original order: 0..22 (23 slides)
    # New slides added at 23, 24, 25, 26
    # Desired order: 0-8 (original 1-9), then 23,24,25,26, then 9-22 (original 10-23)
    new_order = (
        list(range(0, 9)) +  # slides 1–9 (original)
        [23, 24, 25, 26] +  # 4 new progress slides
        list(range(9, 23))      # slides 10–23 (original)
    )
    reorder_slides(prs, new_order)

    # ── 4. Update outline slide (slide 2) to mention Progress Update section ─
    # still index 1 after reorder? Let's check title
    outline_slide = prs.slides[1]
    for shape in outline_slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if "Presentation Outline" in run.text:
                        pass  # title — leave
                    # Inject "Progress Update (V2)" note if outline items found
                    if "Data Collection" in run.text and "✔" not in run.text:
                        run.text += "  ✔"

    DST.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(DST))
    print(f"Saved: {DST}")
    print(f"Total slides: {len(prs.slides)}")
    slide_titles = []
    for s in prs.slides:
        for shape in s.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.paragraphs[0].text.strip()
                if t:
                    slide_titles.append(t[:70])
                    break
        else:
            slide_titles.append("(no text)")
    for i, t in enumerate(slide_titles):
        print(f"  {i+1:2d}. {t}")


if __name__ == "__main__":
    main()
